import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_natural_sih_deck(output_path="SIH2026_Idea_Presentation_PS26154.pptx"):
    prs = Presentation()
    # 16:9 Widescreen (13.333 in x 7.5 in)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Official SIH Minimal Colors
    COLOR_BG = RGBColor(255, 255, 255)           # Pure White
    COLOR_PRIMARY_BLUE = RGBColor(27, 54, 93)    # Official SIH Navy Blue (#1B365D)
    COLOR_ORANGE = RGBColor(243, 112, 33)        # SIH Saffron/Orange (#F37021)
    COLOR_FOOTER_BLUE = RGBColor(0, 119, 200)    # SIH Footer Blue (#0077C8)
    COLOR_TEXT_MAIN = RGBColor(20, 24, 30)       # High contrast dark text (#14181E)

    # EXACT FONTS FROM OFFICIAL TEMPLATE
    FONT_TITLE = "Times New Roman"
    FONT_BODY = "Arial"
    LOGO_PATH = "sih_official_logo.png"
    FLOWCHART_PATH = "omnitransform_pipeline_flowchart.png"

    def apply_header_footer(slide, slide_num, slide_title_text):
        # 1. Top Left Team Oval (Standard, editable, template element)
        team_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(0.3), Inches(1.8), Inches(1.0))
        team_oval.fill.solid()
        team_oval.fill.fore_color.rgb = COLOR_BG
        team_oval.line.color.rgb = COLOR_PRIMARY_BLUE
        team_oval.line.width = Pt(1.5)
        tf_team = team_oval.text_frame
        tf_team.word_wrap = True
        tf_team.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf_team.paragraphs[0]
        p1.text = "Your Team"
        p1.font.size = Pt(11)
        p1.font.name = FONT_BODY
        p1.font.color.rgb = COLOR_TEXT_MAIN
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf_team.add_paragraph()
        p2.text = "Name"
        p2.font.size = Pt(11)
        p2.font.name = FONT_BODY
        p2.font.color.rgb = COLOR_TEXT_MAIN
        p2.alignment = PP_ALIGN.CENTER

        # 2. Top Center Slide Title in Times New Roman (Exact Template Style)
        if slide_title_text:
            tb_title = slide.shapes.add_textbox(Inches(2.5), Inches(0.35), Inches(6.8), Inches(0.85))
            tf_title = tb_title.text_frame
            tf_title.word_wrap = True
            p_t = tf_title.paragraphs[0]
            p_t.text = slide_title_text
            p_t.font.size = Pt(26)
            p_t.font.bold = True
            p_t.font.name = FONT_TITLE
            p_t.font.color.rgb = COLOR_PRIMARY_BLUE
            p_t.alignment = PP_ALIGN.CENTER

        # 3. Top Right Official SIH 2026 Logo (Standalone uploaded logo)
        if os.path.exists(LOGO_PATH):
            slide.shapes.add_picture(LOGO_PATH, Inches(9.8), Inches(0.25), width=Inches(3.0))

        # 4. Bottom Footer Ribbon (Exact SIH Blue #0077C8)
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45))
        footer.fill.solid()
        footer.fill.fore_color.rgb = COLOR_FOOTER_BLUE
        footer.line.fill.background()

        tb_f = slide.shapes.add_textbox(Inches(0.8), Inches(7.08), Inches(10.0), Inches(0.38))
        tf_f = tb_f.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = "@SIH Idea submission- Template"
        p_f.font.size = Pt(10)
        p_f.font.name = FONT_BODY
        p_f.font.color.rgb = RGBColor(255, 255, 255)

        tb_num = slide.shapes.add_textbox(Inches(12.0), Inches(7.08), Inches(0.8), Inches(0.38))
        tf_num = tb_num.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.font.size = Pt(11)
        p_num.font.bold = True
        p_num.font.name = FONT_BODY
        p_num.font.color.rgb = RGBColor(255, 255, 255)
        p_num.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 1: TITLE PAGE (Simple, Clean, Authentic)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)

    tb1_top = slide1.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.6))
    p_t1 = tb1_top.text_frame.paragraphs[0]
    p_t1.text = "SMART INDIA HACKATHON 2026"
    p_t1.font.size = Pt(28)
    p_t1.font.bold = True
    p_t1.font.name = FONT_TITLE
    p_t1.font.color.rgb = COLOR_PRIMARY_BLUE
    p_t1.alignment = PP_ALIGN.CENTER

    tb1_sub = slide1.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.5))
    p_sub = tb1_sub.text_frame.paragraphs[0]
    p_sub.text = "TITLE PAGE"
    p_sub.font.size = Pt(22)
    p_sub.font.bold = True
    p_sub.font.name = FONT_TITLE
    p_sub.font.color.rgb = COLOR_TEXT_MAIN
    p_sub.alignment = PP_ALIGN.CENTER

    # Left: Clean, standard text frame (Easy to click and edit)
    tb1_fields = slide1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7.8), Inches(5.1))
    tf1_f = tb1_fields.text_frame
    tf1_f.word_wrap = True

    s1_items = [
        ("• Problem Statement ID –", " 26154", True),
        ("• Problem Statement Title –", " Gen AI Platform for Automated Content Transformation", True),
        ("• Theme –", " Blockchain & Cybersecurity / AI", False),
        ("• PS Category –", " Software", False),
        ("• Organization –", " National Technical Research Organisation (NTRO)", False),
        ("• Team ID –", " [Registered Team ID on Portal]", False),
        ("• Team Name (Registered on portal) –", " OmniTransform", True),
        ("• Idea Title –", " OmniTransform AI - Automated Multi-Format Content Platform", True)
    ]

    for idx, (label, val, highlight) in enumerate(s1_items):
        p = tf1_f.add_paragraph() if idx > 0 else tf1_f.paragraphs[0]
        r1 = p.add_run()
        r1.text = label
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.name = FONT_TITLE
        r1.font.color.rgb = COLOR_PRIMARY_BLUE

        r2 = p.add_run()
        r2.text = val
        r2.font.size = Pt(14)
        r2.font.bold = highlight
        r2.font.name = FONT_TITLE if highlight else FONT_BODY
        r2.font.color.rgb = COLOR_ORANGE if (highlight and idx == 0) else COLOR_TEXT_MAIN
        p.space_after = Pt(10)

    if os.path.exists(LOGO_PATH):
        slide1.shapes.add_picture(LOGO_PATH, Inches(8.8), Inches(2.0), width=Inches(3.8))

    # =========================================================================
    # SLIDE 2: IDEA TITLE & PROPOSED SOLUTION (Clean, Natural Layout)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_header_footer(slide2, 2, "IDEA TITLE: OmniTransform AI")

    # Main Section Header Pointer (Exact Template Underline)
    tb2_head = slide2.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.5))
    p2_h = tb2_head.text_frame.paragraphs[0]
    p2_h.text = "• Proposed Solution (Describe your Idea/Solution/Prototype)"
    p2_h.font.size = Pt(18)
    p2_h.font.bold = True
    p2_h.font.underline = True
    p2_h.font.name = FONT_TITLE
    p2_h.font.color.rgb = COLOR_PRIMARY_BLUE

    # 3 Clean Columns on plain white background
    col_w = Inches(3.7)
    col_gap = Inches(0.35)
    top_pos = Inches(2.0)
    h_pos = Inches(4.8)

    # Column 1: Detailed Explanation
    tb2_c1 = slide2.shapes.add_textbox(Inches(0.8), top_pos, col_w, h_pos)
    tf2_1 = tb2_c1.text_frame
    tf2_1.word_wrap = True

    p = tf2_1.paragraphs[0]
    p.text = "• Detailed explanation of the proposed solution"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(8)

    p1 = tf2_1.add_paragraph()
    p1.text = "An automated system that ingests long PDFs, security advisories, and policy whitepapers, converting them into 5 ready-to-use formats in under 10 seconds:"
    p1.font.size = Pt(12)
    p1.font.name = FONT_BODY
    p1.font.color.rgb = COLOR_TEXT_MAIN
    p1.space_after = Pt(8)

    outputs = [
        "1. 1-Page Executive Summary Memo",
        "2. Presentation Slide Deck",
        "3. Infographics & Visual Data Cards",
        "4. Multilingual Press Release (English & Regional)",
        "5. 60-Second Voice Audio Briefing"
    ]
    for o in outputs:
        po = tf2_1.add_paragraph()
        po.text = f"   {o}"
        po.font.size = Pt(11.5)
        po.font.name = FONT_BODY
        po.font.color.rgb = COLOR_TEXT_MAIN
        po.space_after = Pt(4)

    # Column 2: How It Addresses Problem
    tb2_c2 = slide2.shapes.add_textbox(Inches(0.8) + col_w + col_gap, top_pos, col_w, h_pos)
    tf2_2 = tb2_c2.text_frame
    tf2_2.word_wrap = True

    p = tf2_2.paragraphs[0]
    p.text = "• How it addresses the problem"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(8)

    pts_address = [
        ("Saves Analyst Time:", "Officers spend 4–6 hours daily converting reports into slides and briefs. OmniTransform automates this in seconds."),
        ("Eliminates Multiple Tools:", "Replaces separate design, translation, and summary tools with one integrated workflow."),
        ("Enables Fast Decision Making:", "Executives and leaders can quickly review key takeaways during critical events."),
        ("Accurate & Verifiable:", "Every generated point links back directly to the source page in the PDF.")
    ]
    for lbl, txt in pts_address:
        p = tf2_2.add_paragraph()
        r1 = p.add_run()
        r1.text = f"- {lbl} "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.name = FONT_BODY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = txt
        r2.font.size = Pt(12)
        r2.font.name = FONT_BODY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(8)

    # Column 3: Innovation and Uniqueness
    tb2_c3 = slide2.shapes.add_textbox(Inches(0.8) + (col_w + col_gap)*2, top_pos, col_w, h_pos)
    tf2_3 = tb2_c3.text_frame
    tf2_3.word_wrap = True

    p = tf2_3.paragraphs[0]
    p.text = "• Innovation and uniqueness of the solution"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(8)

    pts_innov = [
        ("Interactive Source Citations:", "Click any summary point to highlight the exact sentence in the original PDF."),
        ("Audience-Specific Views:", "Switch between an Executive Brief, Citizen Explainer, or Technical Advisory with 1 click."),
        ("Works Completely Offline:", "Can run on local GPU servers without sending sensitive data to external cloud APIs."),
        ("Multi-Format Synchronized Output:", "Generates all 5 formats simultaneously from a single document upload.")
    ]
    for lbl, txt in pts_innov:
        p = tf2_3.add_paragraph()
        r1 = p.add_run()
        r1.text = f"- {lbl} "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.name = FONT_BODY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = txt
        r2.font.size = Pt(12)
        r2.font.name = FONT_BODY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(8)

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH (Clean Flowchart & Real Tech Stack)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_header_footer(slide3, 3, "TECHNICAL APPROACH")

    # Left Container (4.6 in): Technologies to be used
    tb3_l = slide3.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(4.6), Inches(5.5))
    tf3_l = tb3_l.text_frame
    tf3_l.word_wrap = True

    p = tf3_l.paragraphs[0]
    p.text = "• Technologies to be used"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(10)

    techs = [
        ("Frontend:", "Next.js / React, Tailwind CSS, PDF.js document viewer."),
        ("Backend & Parser:", "Python FastAPI, PyMuPDF, pdfplumber for table parsing."),
        ("AI Models:", "Llama 3 (8B/70B) / Mistral, IndicBERT (Indian languages)."),
        ("Slide & Graphic Rendering:", "Marp slide compiler, Canvas HTML-to-Image API."),
        ("Voice Generation:", "Piper Neural TTS (lightweight, runs locally offline)."),
        ("Database & Vector Search:", "FAISS / ChromaDB for document search & citations."),
        ("Deployment:", "Docker containerized, can run air-gapped on Linux/Windows.")
    ]
    for cat, desc in techs:
        p = tf3_l.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{cat} "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.name = FONT_BODY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.name = FONT_BODY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(5)

    # Right Container (7.0 in): Methodology & Clean Flow Chart
    tb3_r_title = slide3.shapes.add_textbox(Inches(5.6), Inches(1.35), Inches(7.0), Inches(0.5))
    p_rt = tb3_r_title.text_frame.paragraphs[0]
    p_rt.text = "• Methodology and process for implementation (Flow Chart)"
    p_rt.font.size = Pt(15)
    p_rt.font.bold = True
    p_rt.font.underline = True
    p_rt.font.name = FONT_TITLE
    p_rt.font.color.rgb = COLOR_PRIMARY_BLUE

    if os.path.exists(FLOWCHART_PATH):
        slide3.shapes.add_picture(FLOWCHART_PATH, Inches(5.6), Inches(1.9), width=Inches(6.9))

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY (Practical Engineering Reality)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_header_footer(slide4, 4, "FEASIBILITY AND VIABILITY")

    # Top Section: Analysis of Feasibility
    tb4_t = slide4.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.7), Inches(2.2))
    tf4_t = tb4_t.text_frame
    tf4_t.word_wrap = True

    p = tf4_t.paragraphs[0]
    p.text = "• Analysis of the feasibility of the idea"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(6)

    feasibility_points = [
        ("Technical Feasibility:", "Uses open-source models and standard document parsers. Tested with sample 30+ page government PDFs."),
        ("Hardware Requirements:", "Runs on a single workstation with an NVIDIA GPU (RTX 3060 or higher) using 4-bit quantized open-source models."),
        ("Deployment Feasibility:", "Packaged with Docker for simple 1-click deployment on internal government networks without cloud setup.")
    ]
    for lbl, desc in feasibility_points:
        p = tf4_t.add_paragraph()
        r1 = p.add_run()
        r1.text = f"- {lbl} "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.name = FONT_BODY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(12)
        r2.font.name = FONT_BODY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(4)

    # Bottom Section: Practical Challenges & Solutions
    tb4_b = slide4.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(3.2))
    tf4_b = tb4_b.text_frame
    tf4_b.word_wrap = True

    p = tf4_b.paragraphs[0]
    p.text = "• Potential challenges and risks & Strategies for overcoming these challenges"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(8)

    risks_matrix = [
        ("1. Risk of AI Hallucinations in Summary Briefs",
         "Challenge: LLMs can sometimes generate inaccurate numbers or dates.",
         "Strategy: Use strict retrieval prompts and bind every output sentence to exact source page citations."),
        
        ("2. Data Privacy & Confidentiality",
         "Challenge: Sensitive official documents cannot be sent to public AI APIs.",
         "Strategy: Host open-source models (Llama 3 / Mistral) entirely on local private servers."),
        
        ("3. Complex Multi-Column Layouts & Tables",
         "Challenge: Standard PDF text extractors scramble tables and multi-column text.",
         "Strategy: Use PyMuPDF spatial coordinate extraction and pdfplumber table bounding boxes.")
    ]

    for title, chall, strat in risks_matrix:
        p = tf4_b.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{title}\n  "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.name = FONT_BODY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        
        r2 = p.add_run()
        r2.text = f"{chall} "
        r2.font.size = Pt(11.5)
        r2.font.name = FONT_BODY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        
        r3 = p.add_run()
        r3.text = f"\n  Solution: {strat}"
        r3.font.bold = True
        r3.font.size = Pt(11.5)
        r3.font.name = FONT_BODY
        r3.font.color.rgb = COLOR_ORANGE
        p.space_after = Pt(6)

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS (Clean, Direct, Practical)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_header_footer(slide5, 5, "IMPACT AND BENEFITS")

    # Left: Impact
    tb5_l = slide5.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(5.7), Inches(5.4))
    tf5_l = tb5_l.text_frame
    tf5_l.word_wrap = True

    p = tf5_l.paragraphs[0]
    p.text = "• Potential impact on the target audience"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(10)

    auds = [
        ("Defense & Intelligence Analysts (NTRO / MoD):", "Saves 4 to 6 hours daily per analyst by instantly synthesizing long threat advisories and technical reports."),
        ("Senior Officials & Executives:", "Provides 1-page summaries and ready-to-present slide decks for urgent decision meetings."),
        ("Public Relations & Media Officers:", "Quickly converts technical policy documents into public press releases and infographics in English and regional languages."),
        ("Field Staff & Mobile Teams:", "Allows personnel to listen to 60-second audio podcast briefings on the go.")
    ]
    for aud, text in auds:
        p = tf5_l.add_paragraph()
        r1 = p.add_run()
        r1.text = f"- {aud}\n  "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.name = FONT_BODY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(11.5)
        r2.font.name = FONT_BODY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(6)

    # Right: Benefits
    tb5_r = slide5.shapes.add_textbox(Inches(6.8), Inches(1.35), Inches(5.7), Inches(5.4))
    tf5_r = tb5_r.text_frame
    tf5_r.word_wrap = True

    p = tf5_r.paragraphs[0]
    p.text = "• Benefits of the solution (social, economic, etc.)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(10)

    bens = [
        ("Time & Cost Savings:", "Reduces manual document formatting and translation time by over 80%, saving departmental resources."),
        ("Regional Language Accessibility:", "Bridges language barriers by translating government policies into Hindi, Kannada, and other Indian languages."),
        ("Enhanced Information Security:", "100% data remains on internal servers with zero risk of leaks to third-party APIs."),
        ("Eco-Friendly & Paperless:", "Reduces printing of bulky multi-page reports by providing interactive digital memos and presentations.")
    ]
    for lbl, text in bens:
        p = tf5_r.add_paragraph()
        r1 = p.add_run()
        r1.text = f"- {lbl}\n  "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.name = FONT_BODY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(11.5)
        r2.font.name = FONT_BODY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(6)

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES (Real Research & Exact Links)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_header_footer(slide6, 6, "RESEARCH AND REFERENCES")

    # Left: Research Papers
    tb6_l = slide6.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(5.7), Inches(5.4))
    tf6_l = tb6_l.text_frame
    tf6_l.word_wrap = True

    p = tf6_l.paragraphs[0]
    p.text = "• Details / Links of the reference and research work"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(10)

    research_citations = [
        ("Retrieval-Augmented Generation (RAG):", "Lewis et al. (NeurIPS 2020) — Foundation for verifiable knowledge retrieval."),
        ("Open-Weight Foundation Models:", "Meta Llama 3 & Mistral AI Research Papers on Instruction Tuning (2024)."),
        ("Multilingual NLP for Indian Languages:", "AI4Bharat & Bhashini Project Research on IndicBERT (2023)."),
        ("Document Layout Analysis:", "PyMuPDF spatial coordinate mapping and bounding box extraction standards.")
    ]
    for cit, title in research_citations:
        p = tf6_l.add_paragraph()
        r1 = p.add_run()
        r1.text = f"- {cit}\n  "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.name = FONT_BODY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = title
        r2.font.size = Pt(11.5)
        r2.font.name = FONT_BODY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(6)

    # Right: Project Links
    tb6_r = slide6.shapes.add_textbox(Inches(6.8), Inches(1.35), Inches(5.7), Inches(5.4))
    tf6_r = tb6_r.text_frame
    tf6_r.word_wrap = True

    p = tf6_r.paragraphs[0]
    p.text = "• Official Portals & Project Verification Links"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_TITLE
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(10)

    links_data = [
        ("SIH 2026 Problem Statement 26154:", "https://www.sih.gov.in/sih2026PS"),
        ("NTRO Official Government Website:", "https://ntro.gov.in"),
        ("Live Web Prototype & Design:", "https://hrl-brand-seo.vercel.app"),
        ("GitHub Project Repository:", "https://github.com/hrlpavan/sih-2026-media-problem-statements"),
        ("Winning Strategy & Technical Documentation:", "https://github.com/hrlpavan/sih-2026-media-problem-statements/blob/main/WINNING_STRATEGY_PS26154.md")
    ]
    for lbl, url in links_data:
        p = tf6_r.add_paragraph()
        r1 = p.add_run()
        r1.text = f"- {lbl}\n  "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.name = FONT_BODY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = url
        r2.font.size = Pt(11)
        r2.font.name = FONT_BODY
        r2.font.color.rgb = COLOR_FOOTER_BLUE
        p.space_after = Pt(6)

    prs.save(output_path)
    print(f"Natural, clean SIH 2026 deck created successfully at: {output_path}")

if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "SIH2026_Idea_Presentation_PS26154.pptx"
    build_natural_sih_deck(out)
