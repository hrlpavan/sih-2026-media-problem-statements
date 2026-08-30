import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck(output_path="OmniTransform_AI_Pitch_Deck.pptx"):
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette: Deep Blue / Dark Modern Theme
    BG_DARK = RGBColor(10, 15, 29)        # #0A0F1D
    CARD_DARK = RGBColor(19, 27, 49)      # #131B31
    ACCENT_BLUE = RGBColor(0, 113, 227)   # #0071E3 (Apple Blue)
    ACCENT_CYAN = RGBColor(0, 210, 255)   # #00D2FF
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(160, 175, 200)
    ACCENT_GOLD = RGBColor(255, 184, 0)

    slides_data = [
        {
            "badge": "SIH 2026 / EUREKA! 2026 PITCH DECK",
            "title": "OmniTransform AI",
            "subtitle": "Sovereign Multimodal Content Transformation Platform for Government & Enterprise Intelligence",
            "points": [
                "Target Problem Statement: PS ID 26154 (NTRO)",
                "Category: Software | Theme: Blockchain & Cybersecurity / AI",
                "Team Leader: Pavan Kumar S | SCEM / Sahyadri Ecosystem",
                "Live Platform Ecosystem: https://hrl-brand-seo.vercel.app"
            ],
            "is_title": True
        },
        {
            "badge": "THE PROBLEM & BOTTLENECK",
            "title": "Dense Documents, Slow Manual Synthesis",
            "subtitle": "High-level analysts and officers waste 4-6 hours daily translating documents into separate formats.",
            "points": [
                "4-6 Hours Wasted Daily: Translating raw 50+ page PDFs into briefings, decks, and press releases.",
                "Fragmented Tooling: Juggling Canva, PowerPoint, Word, and Chatbots leads to inconsistent messaging.",
                "Tactical Information Decay: Critical threat advisories and policy whitepapers lose impact due to delayed delivery.",
                "The Hallucination Danger: Public AI tools invent facts, making unverified outputs dangerous for governance."
            ]
        },
        {
            "badge": "THE SOLUTION",
            "title": "OmniTransform AI: 1 Ingestion  5 Synchronized Assets",
            "subtitle": "Drag and drop any 30+ page PDF; receive 5 production-ready media formats in under 10 seconds.",
            "points": [
                "1.  1-Page Executive Memo: High-level tactical summary with actionable bullet points.",
                "2.  Interactive Presentation Deck: Keynote-grade slide layout ready for leadership meetings.",
                "3.  Visual Infographic & Social Card: Auto-rendered canvas assets for public and press outreach.",
                "4.  Multilingual Press Release: Automated regional translation (English, Hindi, Indic languages).",
                "5.  60-Second AI Voice Intelligence Briefing: Studio-quality audio podcast for mobile officials."
            ]
        },
        {
            "badge": "TECHNICAL ARCHITECTURE",
            "title": "End-to-End Multimodal Ingestion Pipeline",
            "subtitle": "Modular, sovereign, and air-gapped architecture with zero cloud dependencies.",
            "points": [
                "Ingestion Layer: PyMuPDF & Unstructured parser with exact bounding-box coordinate tracking.",
                "Reasoning & LLM Pipeline: LangChain / LlamaIndex with constrained JSON schema enforcement.",
                "Dynamic Slide & Canvas Compilers: Satori & Marp engines for instant vector image/slide rendering.",
                "Voice AI Synthesis Engine: Neural TTS & ElevenLabs API for rapid studio-grade audio synthesis.",
                "Reverse Coordinate Verification: Interactive PDF highlighter providing 100% citation grounding."
            ]
        },
        {
            "badge": "KEY INNOVATION & USPs",
            "title": "Why OmniTransform AI is Unbeatable",
            "subtitle": "Engineered to eliminate judge skepticism and solve core enterprise pain points.",
            "points": [
                "Zero-Hallucination Source Grounding: Click any insight badge to highlight the exact original PDF sentence.",
                "Multi-Persona Tone Switcher: 1-click toggle between Executive, Technical, Citizen, and Press formats.",
                "10-Second Deterministic Runtime: Instant local rendering eliminates venue Wi-Fi failure during live demos.",
                "Apple-Grade Human Interface: Glassmorphism aesthetics, responsive drawer navigation, and high contrast."
            ]
        },
        {
            "badge": "TECHNOLOGY STACK",
            "title": "Production-Grade Modern Tech Stack",
            "subtitle": "Robust, containerized, and optimized for speed and on-premise security.",
            "points": [
                "Frontend: Next.js 14, React, Tailwind CSS, Lucide Icons, PDF.js canvas viewer.",
                "Backend API: Python FastAPI, LangChain, PyMuPDF, Qdrant / FAISS Vector Store.",
                "AI Models: Llama 3.3 (70B), Mistral Large, IndicBERT for regional Indic language processing.",
                "Audio & Media: ElevenLabs / Piper Neural TTS, FFmpeg, Canvas HTML-to-Image API.",
                "Deployment: Docker, Sovereign Air-Gapped Linux Host, Zero External Telemetry."
            ]
        },
        {
            "badge": "FEASIBILITY & RISK MITIGATION",
            "title": "Proactive Engineering & Risk Defense",
            "subtitle": "Addressing enterprise reliability, data security, and performance constraints.",
            "points": [
                "Risk: AI Hallucination  Solution: Strict RAG boundaries and regex coordinate verification.",
                "Risk: Sensitive Data Leaks  Solution: 100% on-premise air-gapped open-weight model deployment.",
                "Risk: Complex Charts & Tables  Solution: Multimodal Vision-Language Models (VLM) for diagram extraction.",
                "Risk: Large File Processing Latency  Solution: Map-reduce parallelized document chunking pipeline."
            ]
        },
        {
            "badge": "BUSINESS & MARKET POTENTIAL",
            "title": "B2G / B2B Market Scalability",
            "subtitle": "Massive $109B global market with high-urgency government and enterprise demand.",
            "points": [
                "Market Size: $109B Global Enterprise Content & GenAI market by 2030 ($4.2B APAC SAM).",
                "Target Sectors: Defense/Intelligence (NTRO, MoD, I4C), Corporate Consulting, BFSI, and Media Houses.",
                "B2G License Model: ₹15L - ₹45L/year for on-premise air-gapped sovereign deployment.",
                "B2B SaaS Model: $49 - $299/user/month for consulting firms, research desks, and legal teams.",
                "Traction Projection: Leveraging HRL International's infrastructure to achieve ₹3.5 Cr ARR by Year 2."
            ]
        },
        {
            "badge": "EXECUTION & TEAM ROLES",
            "title": "Balanced 6-Member Team Execution Plan",
            "subtitle": "Clear division of engineering, design, and presentation responsibilities.",
            "points": [
                "Pavan Kumar S (Team Lead): System Architecture, Pitch Delivery, Apple UI/UX Integration.",
                "Frontend Engineer: Next.js 14 Dashboard, File Dropzone, Multi-Asset Preview Tabs.",
                "Backend Engineer: FastAPI Server, PyMuPDF Document Coordinate Parser & API Routes.",
                "AI/LLM Engineer: LangChain Extraction Pipeline, Structured JSON Schema Prompts.",
                "Visual & Canvas Engineer: Satori Infographic Auto-Renderer & Slide Layout Engine.",
                "Voice & Translation Lead: ElevenLabs/Piper TTS Integration & Multilingual Regional NLP."
            ]
        },
        {
            "badge": "CONCLUSION & VERIFICATION",
            "title": "Why OmniTransform AI Wins SIH / Eureka! 2026",
            "subtitle": "A verified, live, and battle-tested ecosystem ready for real-world deployment.",
            "points": [
                "Solves a Critical National Need: Streamlines intelligence dissemination for NTRO & government bodies.",
                "10-Second Visual WOW Factor: Generates 5 high-impact media formats live in front of the jury.",
                "Zero Hallucination Proof: Verified source citations directly linking to raw document text.",
                "Live Platform: https://hrl-brand-seo.vercel.app | Repo: https://github.com/hrlpavan/sih-2026-media-problem-statements"
            ]
        }
    ]

    for data in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # 1. Background rectangle
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()

        # 2. Top Header Accent Pill / Badge
        badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf_badge = badge_box.text_frame
        tf_badge.word_wrap = True
        p_badge = tf_badge.paragraphs[0]
        p_badge.text = data["badge"].upper()
        p_badge.font.size = Pt(11)
        p_badge.font.bold = True
        p_badge.font.color.rgb = ACCENT_CYAN
        p_badge.font.name = "Arial"

        # 3. Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.9))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = data["title"]
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        p_title.font.name = "Arial"

        # 4. Slide Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.65), Inches(11.7), Inches(0.6))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = data["subtitle"]
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = TEXT_MUTED
        p_sub.font.name = "Arial"

        # 5. Main Card Container
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.35), Inches(11.733), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_DARK
        card.line.color.rgb = RGBColor(35, 48, 82)
        card.line.width = Pt(1.5)

        # 6. Bullet Points / Content Inside Card
        content_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.55), Inches(11.1), Inches(4.1))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True

        for i, pt in enumerate(data["points"]):
            p = tf_content.add_paragraph() if i > 0 else tf_content.paragraphs[0]
            p.text = pt
            p.font.size = Pt(15 if not data.get("is_title") else 17)
            p.font.color.rgb = TEXT_WHITE
            p.font.name = "Arial"
            p.space_after = Pt(14)
            p.level = 0

    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "OmniTransform_AI_Pitch_Deck.pptx"
    create_deck(out)
