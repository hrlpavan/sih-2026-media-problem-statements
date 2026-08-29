import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_golden_ratio_sih_deck(output_path="SIH2026_Idea_Presentation_PS26154.pptx"):
    prs = Presentation()
    # 16:9 Widescreen (13.333 in x 7.5 in)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # STRICT COLOR PALETTE (Extracted directly from the official SIH 2026 Template)
    # No foreign/unwanted colors.
    COLOR_BG = RGBColor(255, 255, 255)           # Pure White (#FFFFFF)
    COLOR_PRIMARY_BLUE = RGBColor(27, 54, 93)    # Official SIH Header Blue (#1B365D)
    COLOR_ORANGE = RGBColor(243, 112, 33)        # SIH Saffron/Orange Accent (#F37021)
    COLOR_GREEN = RGBColor(30, 168, 88)          # SIH Green (#1EA858)
    COLOR_FOOTER_BLUE = RGBColor(0, 119, 200)    # SIH Footer Bar Blue (#0077C8)
    COLOR_TEXT_MAIN = RGBColor(33, 37, 41)       # Charcoal Black for High Legibility (#212529)
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139)   # Subdued Slate for Labels (#64748B)
    COLOR_CARD_BG = RGBColor(248, 250, 252)      # Minimal Clean Card Fill (#F8FAFC)
    COLOR_CARD_BORDER = RGBColor(226, 232, 240)  # Subtle Slate Border (#E2E8F0)

    # GOLDEN RATIO TYPOGRAPHY SCALE (phi = 1.618)
    # Base: 11pt -> 14pt (x1.27) -> 18pt (x1.618) -> 26pt (x1.618^2)
    FONT_TITLE = Pt(26)
    FONT_SECTION_HEADER = Pt(15.5)
    FONT_SUB_HEADER = Pt(13)
    FONT_BODY = Pt(11)
    FONT_SMALL = Pt(9.5)
    FONT_FAMILY = "Arial"

    def apply_header_footer(slide, slide_num, slide_title_text):
        # 1. Top Left Oval: "Your Team Name" (Exactly as in official template)
        team_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), Inches(0.35), Inches(1.75), Inches(1.05))
        team_oval.fill.solid()
        team_oval.fill.fore_color.rgb = COLOR_BG
        team_oval.line.color.rgb = COLOR_PRIMARY_BLUE
        team_oval.line.width = Pt(1.5)
        tf_team = team_oval.text_frame
        tf_team.word_wrap = True
        tf_team.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf_team.paragraphs[0]
        p1.text = "Team"
        p1.font.size = Pt(10)
        p1.font.name = FONT_FAMILY
        p1.font.color.rgb = COLOR_TEXT_MAIN
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf_team.add_paragraph()
        p2.text = "OmniTransform"
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.name = FONT_FAMILY
        p2.font.color.rgb = COLOR_PRIMARY_BLUE
        p2.alignment = PP_ALIGN.CENTER

        # 2. Top Center Slide Title
        if slide_title_text:
            tb_title = slide.shapes.add_textbox(Inches(2.5), Inches(0.45), Inches(7.8), Inches(0.8))
            tf_title = tb_title.text_frame
            tf_title.word_wrap = True
            p_t = tf_title.paragraphs[0]
            p_t.text = slide_title_text
            p_t.font.size = Pt(24)
            p_t.font.bold = True
            p_t.font.name = FONT_FAMILY
            p_t.font.color.rgb = COLOR_PRIMARY_BLUE
            p_t.alignment = PP_ALIGN.CENTER

        # 3. Top Right SIH Logo & Text (Exactly as in official template)
        tb_sih = slide.shapes.add_textbox(Inches(10.5), Inches(0.3), Inches(2.4), Inches(0.9))
        tf_sih = tb_sih.text_frame
        tf_sih.word_wrap = True
        p_s1 = tf_sih.paragraphs[0]
        p_s1.text = "SMART INDIA"
        p_s1.font.size = Pt(10)
        p_s1.font.bold = True
        p_s1.font.name = FONT_FAMILY
        p_s1.font.color.rgb = COLOR_PRIMARY_BLUE
        p_s1.alignment = PP_ALIGN.RIGHT
        p_s2 = tf_sih.add_paragraph()
        p_s2.text = "HACKATHON"
        p_s2.font.size = Pt(10)
        p_s2.font.bold = True
        p_s2.font.name = FONT_FAMILY
        p_s2.font.color.rgb = COLOR_PRIMARY_BLUE
        p_s2.alignment = PP_ALIGN.RIGHT
        p_s3 = tf_sih.add_paragraph()
        p_s3.text = "2026"
        p_s3.font.size = Pt(12)
        p_s3.font.bold = True
        p_s3.font.name = FONT_FAMILY
        p_s3.font.color.rgb = COLOR_ORANGE
        p_s3.alignment = PP_ALIGN.RIGHT

        # 4. Bottom Footer Ribbon (Exact SIH Blue #0077C8)
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45))
        footer.fill.solid()
        footer.fill.fore_color.rgb = COLOR_FOOTER_BLUE
        footer.line.fill.background()

        tb_f = slide.shapes.add_textbox(Inches(0.8), Inches(7.08), Inches(10.0), Inches(0.38))
        tf_f = tb_f.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = "@SIH Idea submission- Template"
        p_f.font.size = Pt(9.5)
        p_f.font.name = FONT_FAMILY
        p_f.font.color.rgb = RGBColor(255, 255, 255)

        tb_num = slide.shapes.add_textbox(Inches(12.0), Inches(7.08), Inches(0.8), Inches(0.38))
        tf_num = tb_num.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.font.size = Pt(10.5)
        p_num.font.bold = True
        p_num.font.name = FONT_FAMILY
        p_num.font.color.rgb = RGBColor(255, 255, 255)
        p_num.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 1: TITLE PAGE (Rule of Thirds 2-Column Balanced Grid)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)

    # Top Header
    tb1_top = slide1.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(10.0), Inches(0.6))
    p_t1 = tb1_top.text_frame.paragraphs[0]
    p_t1.text = "SMART INDIA HACKATHON 2026"
    p_t1.font.size = Pt(28)
    p_t1.font.bold = True
    p_t1.font.name = FONT_FAMILY
    p_t1.font.color.rgb = COLOR_PRIMARY_BLUE

    # Top Right SIH Text
    tb1_r = slide1.shapes.add_textbox(Inches(10.5), Inches(0.35), Inches(2.4), Inches(0.8))
    p_tr1 = tb1_r.text_frame.paragraphs[0]
    p_tr1.text = "SMART INDIA\nHACKATHON 2026"
    p_tr1.font.size = Pt(10)
    p_tr1.font.bold = True
    p_tr1.font.name = FONT_FAMILY
    p_tr1.font.color.rgb = COLOR_PRIMARY_BLUE
    p_tr1.alignment = PP_ALIGN.RIGHT

    # Subtitle Center
    tb1_sub = slide1.shapes.add_textbox(Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.5))
    p_sub = tb1_sub.text_frame.paragraphs[0]
    p_sub.text = "TITLE PAGE"
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.name = FONT_FAMILY
    p_sub.font.color.rgb = COLOR_TEXT_MAIN
    p_sub.alignment = PP_ALIGN.CENTER

    # Left Column (2/3 width = 7.6 inches): Mandatory Official Fields
    card_s1_left = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.88), Inches(1.8), Inches(7.6), Inches(5.0))
    card_s1_left.fill.solid()
    card_s1_left.fill.fore_color.rgb = COLOR_CARD_BG
    card_s1_left.line.color.rgb = COLOR_CARD_BORDER
    card_s1_left.line.width = Pt(1.2)

    tf_s1 = card_s1_left.text_frame
    tf_s1.word_wrap = True
    tf_s1.margin_left = Inches(0.35)
    tf_s1.margin_top = Inches(0.3)
    tf_s1.margin_right = Inches(0.35)

    s1_fields = [
        ("• Problem Statement ID –", " 26154", True),
        ("• Problem Statement Title –", " Gen AI Platform for Automated Content Transformation", True),
        ("• Theme –", " Blockchain & Cybersecurity / AI", False),
        ("• PS Category –", " Software", False),
        ("• Organization –", " National Technical Research Organisation (NTRO)", False),
        ("• Team ID –", " [Registered Team ID]", False),
        ("• Team Name (Registered on portal) –", " OmniTransform", True)
    ]

    for idx, (label, val, highlight) in enumerate(s1_fields):
        p = tf_s1.add_paragraph() if idx > 0 else tf_s1.paragraphs[0]
        r1 = p.add_run()
        r1.text = label
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.name = FONT_FAMILY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE

        r2 = p.add_run()
        r2.text = val
        r2.font.size = Pt(12)
        r2.font.bold = highlight
        r2.font.name = FONT_FAMILY
        r2.font.color.rgb = COLOR_ORANGE if highlight and idx == 0 else COLOR_TEXT_MAIN
        p.space_after = Pt(10)

    # Right Column (1/3 width = 3.6 inches): Solution Emblem Card
    card_s1_r = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.75), Inches(1.8), Inches(3.7), Inches(5.0))
    card_s1_r.fill.solid()
    card_s1_r.fill.fore_color.rgb = COLOR_PRIMARY_BLUE
    card_s1_r.line.fill.background()

    tf_s1_r = card_s1_r.text_frame
    tf_s1_r.word_wrap = True
    tf_s1_r.margin_left = Inches(0.25)
    tf_s1_r.margin_top = Inches(0.4)
    tf_s1_r.margin_right = Inches(0.25)

    p_rt = tf_s1_r.paragraphs[0]
    p_rt.text = "PROPOSED SYSTEM"
    p_rt.font.size = Pt(11)
    p_rt.font.bold = True
    p_rt.font.name = FONT_FAMILY
    p_rt.font.color.rgb = COLOR_ORANGE
    p_rt.alignment = PP_ALIGN.CENTER
    p_rt.space_after = Pt(6)

    p_rn = tf_s1_r.add_paragraph()
    p_rn.text = "OmniTransform AI"
    p_rn.font.size = Pt(19)
    p_rn.font.bold = True
    p_rn.font.name = FONT_FAMILY
    p_rn.font.color.rgb = RGBColor(255, 255, 255)
    p_rn.alignment = PP_ALIGN.CENTER
    p_rn.space_after = Pt(14)

    s1_key_points = [
        "1-to-5 Synchronized Ingestion Engine",
        "100% Grounded Source Citations",
        "Air-Gapped Sovereign On-Premise",
        "Live Proof: hrl-brand-seo.vercel.app"
    ]
    for kp in s1_key_points:
        p = tf_s1_r.add_paragraph()
        p.text = f"✔ {kp}"
        p.font.size = Pt(10.5)
        p.font.name = FONT_FAMILY
        p.font.color.rgb = RGBColor(225, 238, 255)
        p.space_after = Pt(12)

    # =========================================================================
    # SLIDE 2: PROPOSED SOLUTION (Rule of Thirds 3-Column Layout)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_header_footer(slide2, 2, "IDEA TITLE")

    # Section Subheader (Exact from Template)
    tb2_sub = slide2.shapes.add_textbox(Inches(0.88), Inches(1.3), Inches(11.57), Inches(0.45))
    p2_head = tb2_sub.text_frame.paragraphs[0]
    p2_head.text = "• Proposed Solution (Describe your Idea/Solution/Prototype)"
    p2_head.font.size = FONT_SECTION_HEADER
    p2_head.font.bold = True
    p2_head.font.underline = True
    p2_head.font.name = FONT_FAMILY
    p2_head.font.color.rgb = COLOR_PRIMARY_BLUE

    # 3 Golden Ratio Columns (Width = 3.65 in, Gap = 0.31 in)
    c_w = Inches(3.65)
    c_gap = Inches(0.31)
    s2_left = Inches(0.88)
    s2_top = Inches(1.85)
    s2_h = Inches(4.95)

    # Column 1: Detailed Explanation
    c1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s2_left, s2_top, c_w, s2_h)
    c1.fill.solid()
    c1.fill.fore_color.rgb = COLOR_CARD_BG
    c1.line.color.rgb = COLOR_CARD_BORDER
    tf_c1 = c1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_left = Inches(0.2)
    tf_c1.margin_top = Inches(0.2)
    tf_c1.margin_right = Inches(0.2)

    p = tf_c1.paragraphs[0]
    p.text = "Detailed Explanation"
    p.font.size = FONT_SUB_HEADER
    p.font.bold = True
    p.font.name = FONT_FAMILY
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(8)

    pts1 = [
        ("Core Function:", " Ingests raw 50+ page PDFs, threat advisories, and policy documents in a single pass."),
        ("1-to-5 Transformation:", " Automatically produces 5 tailored communication assets in <10s:"),
        ("1.", " 📊 Executive 1-Page Summary Memo"),
        ("2.", " 🖥️ Meeting-Ready Slide Deck"),
        ("3.", " 🎨 Public Visual Infographic Cards"),
        ("4.", " 📰 Regional Multilingual Press Release"),
        ("5.", " 🎙️ 60s Voice AI Audio Briefing")
    ]
    for idx, (label, text) in enumerate(pts1):
        p = tf_c1.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{label} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = FONT_FAMILY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE if idx < 2 else COLOR_ORANGE
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(10)
        r2.font.name = FONT_FAMILY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(4)

    # Column 2: How It Addresses Problem
    c2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s2_left + c_w + c_gap, s2_top, c_w, s2_h)
    c2.fill.solid()
    c2.fill.fore_color.rgb = COLOR_CARD_BG
    c2.line.color.rgb = COLOR_CARD_BORDER
    tf_c2 = c2.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_left = Inches(0.2)
    tf_c2.margin_top = Inches(0.2)
    tf_c2.margin_right = Inches(0.2)

    p = tf_c2.paragraphs[0]
    p.text = "How It Addresses Problem"
    p.font.size = FONT_SUB_HEADER
    p.font.bold = True
    p.font.name = FONT_FAMILY
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(8)

    pts2 = [
        ("Eliminates Manual Bottleneck:", " Saves 4-6 hours daily spent by officers rewriting dense reports for leadership and public channels."),
        ("Unified Sovereign Dashboard:", " Replaces disconnected tools (Canva, Word, PPT, Chatbots) into one unified system."),
        ("Prevents Tactical Delay:", " Delivers instant threat intelligence dissemination for time-critical defense scenarios."),
        ("Eliminates Hallucinations:", " Strict RAG boundaries guarantee outputs are 100% faithful to the source document.")
    ]
    for label, text in pts2:
        p = tf_c2.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {label} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = FONT_FAMILY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(10)
        r2.font.name = FONT_FAMILY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(6)

    # Column 3: Innovation and Uniqueness
    c3 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s2_left + (c_w + c_gap)*2, s2_top, c_w, s2_h)
    c3.fill.solid()
    c3.fill.fore_color.rgb = COLOR_CARD_BG
    c3.line.color.rgb = COLOR_CARD_BORDER
    tf_c3 = c3.text_frame
    tf_c3.word_wrap = True
    tf_c3.margin_left = Inches(0.2)
    tf_c3.margin_top = Inches(0.2)
    tf_c3.margin_right = Inches(0.2)

    p = tf_c3.paragraphs[0]
    p.text = "Innovation & Uniqueness"
    p.font.size = FONT_SUB_HEADER
    p.font.bold = True
    p.font.name = FONT_FAMILY
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(8)

    pts3 = [
        ("Interactive Source Highlighting:", " Click any bullet point to open the PDF with exact sentence coordinates highlighted."),
        ("Dynamic Multi-Persona Switching:", " 1-click toggle between C-Suite Brief, Citizen Explainer, and Threat Alert."),
        ("Air-Gapped Sovereign Security:", " Operates entirely offline on open-weights (Llama 3/Mistral) with zero cloud telemetry."),
        ("Apple Design Standard:", " High-contrast glassmorphism UI engineered for high readability and fast decision-making.")
    ]
    for label, text in pts3:
        p = tf_c3.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {label} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = FONT_FAMILY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(10)
        r2.font.name = FONT_FAMILY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(6)

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH (Golden Ratio 40/60 Split Grid)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_header_footer(slide3, 3, "TECHNICAL APPROACH")

    # Left Container (4.6 inches = ~40%): Technologies Used
    c3_l = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.88), Inches(1.45), Inches(4.7), Inches(5.35))
    c3_l.fill.solid()
    c3_l.fill.fore_color.rgb = COLOR_CARD_BG
    c3_l.line.color.rgb = COLOR_CARD_BORDER
    tf3_l = c3_l.text_frame
    tf3_l.word_wrap = True
    tf3_l.margin_left = Inches(0.25)
    tf3_l.margin_top = Inches(0.2)
    tf3_l.margin_right = Inches(0.25)

    p = tf3_l.paragraphs[0]
    p.text = "• Technologies to be used"
    p.font.size = FONT_SECTION_HEADER
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_FAMILY
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(10)

    techs = [
        ("Frontend:", "Next.js 14, React, Tailwind CSS, Lucide Icons, PDF.js canvas viewer."),
        ("Backend & Parser:", "Python FastAPI, PyMuPDF spatial coordinate extractor, LangChain / LlamaIndex."),
        ("AI / Foundation LLMs:", "Llama 3.3 (70B), Mistral Large, IndicBERT for regional translation."),
        ("Slide & Canvas Renderers:", "Satori HTML-to-Image Canvas API, Marp slide compiler, FFmpeg."),
        ("Voice AI Engine:", "Piper Neural TTS & ElevenLabs API."),
        ("Infrastructure:", "Docker containerized, air-gapped on-premise Linux host, zero telemetry.")
    ]
    for lbl, desc in techs:
        p = tf3_l.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{lbl} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = FONT_FAMILY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(9.5)
        r2.font.name = FONT_FAMILY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(5)

    # Right Container (6.6 inches = ~60%): Process Flow & Methodology
    c3_r = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.85), Inches(1.45), Inches(6.6), Inches(5.35))
    c3_r.fill.solid()
    c3_r.fill.fore_color.rgb = COLOR_CARD_BG
    c3_r.line.color.rgb = COLOR_CARD_BORDER
    tf3_r = c3_r.text_frame
    tf3_r.word_wrap = True
    tf3_r.margin_left = Inches(0.25)
    tf3_r.margin_top = Inches(0.2)
    tf3_r.margin_right = Inches(0.25)

    p = tf3_r.paragraphs[0]
    p.text = "• Methodology & Implementation Process Flow"
    p.font.size = FONT_SECTION_HEADER
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_FAMILY
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(10)

    steps = [
        ("Step 1: Ingestion & Spatial Tagging", "PyMuPDF parses raw PDF/DOCX, tagging text blocks, tables, and diagrams with exact coordinate bounding boxes."),
        ("Step 2: Hierarchical Chunking & RAG", "Contextual embeddings are indexed into an in-memory vector store (FAISS/Qdrant) enforcing strict citation boundaries."),
        ("Step 3: Multi-Persona Orchestration", "Constrained JSON Schema prompts instruct the LLM to extract key insights, slide structures, and press summaries."),
        ("Step 4: Parallel Synthesis Engines", "Simultaneously compiles HTML5 presentation slides, auto-renders visual cards (Canvas API), and synthesizes neural audio."),
        ("Step 5: Interactive Verification UI", "Delivers a live dashboard with side-by-side asset previews and clickable citation coordinates directly linking to the source PDF.")
    ]
    for st, desc in steps:
        p = tf3_r.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {st}: "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = FONT_FAMILY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(9.5)
        r2.font.name = FONT_FAMILY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(6)

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY (Rule of Thirds Hierarchy)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_header_footer(slide4, 4, "FEASIBILITY AND VIABILITY")

    # Section 1: Feasibility Analysis (Top Half)
    c4_t = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.88), Inches(1.4), Inches(11.57), Inches(2.2))
    c4_t.fill.solid()
    c4_t.fill.fore_color.rgb = COLOR_CARD_BG
    c4_t.line.color.rgb = COLOR_CARD_BORDER
    tf4_t = c4_t.text_frame
    tf4_t.word_wrap = True
    tf4_t.margin_left = Inches(0.25)
    tf4_t.margin_top = Inches(0.18)
    tf4_t.margin_right = Inches(0.25)

    p = tf4_t.paragraphs[0]
    p.text = "• Analysis of the feasibility of the idea"
    p.font.size = FONT_SECTION_HEADER
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_FAMILY
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(6)

    feas_points = [
        ("Hardware Feasibility:", " Runs efficiently on commodity GPU workstations (e.g. Single RTX 4090 or Apple Silicon) using 4-bit quantized open-weight models."),
        ("Speed & Throughput:", " Asynchronous map-reduce pipeline completes multi-format transformation of a 30-page PDF in < 10 seconds."),
        ("Operational Feasibility:", " 100% Dockerized container stack; operates fully air-gapped without external internet dependencies.")
    ]
    for lbl, desc in feas_points:
        p = tf4_t.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {lbl} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = FONT_FAMILY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10)
        r2.font.name = FONT_FAMILY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(4)

    # Section 2: Challenges & Mitigation Strategies (Bottom Half - 3 Cards)
    c4_b_w = Inches(3.65)
    c4_b_gap = Inches(0.31)
    s4_b_top = Inches(3.8)
    s4_b_h = Inches(3.0)

    risks = [
        ("Potential Risk: AI Hallucination", "Risk: Generative models inventing ungrounded facts in security briefs.", "Mitigation Strategy: Strict RAG constraint boundaries + reverse citation bounding-box verification for 100% grounded truth."),
        ("Potential Risk: Data Security", "Risk: Sensitive intelligence leaked to public cloud APIs.", "Mitigation Strategy: Sovereign on-premise air-gapped deployment with local open-weights (Llama 3/Mistral); zero telemetry."),
        ("Potential Risk: Complex Tables/Charts", "Risk: Loss of tabular structure and diagram metrics during parsing.", "Mitigation Strategy: Multimodal vision-language parsing (VLM) + Markdown table structure extraction.")
    ]

    for i, (r_title, r_chall, r_strat) in enumerate(risks):
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, s2_left + (c4_b_w + c4_b_gap)*i, s4_b_top, c4_b_w, s4_b_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.18)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = r_title
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.name = FONT_FAMILY
        p.font.color.rgb = COLOR_ORANGE
        p.space_after = Pt(6)

        p1 = tf.add_paragraph()
        p1.text = f"• {r_chall}"
        p1.font.size = Pt(9.5)
        p1.font.name = FONT_FAMILY
        p1.font.color.rgb = COLOR_TEXT_MAIN
        p1.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = f"• {r_strat}"
        p2.font.size = Pt(9.5)
        p2.font.bold = True
        p2.font.name = FONT_FAMILY
        p2.font.color.rgb = COLOR_PRIMARY_BLUE

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS (Golden Ratio 50/50 Dual Column)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_header_footer(slide5, 5, "IMPACT AND BENEFITS")

    col_5_w = Inches(5.63)
    col_5_gap = Inches(0.31)

    # Left Container: Potential Impact on Target Audience
    c5_l = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.88), Inches(1.4), col_5_w, Inches(5.4))
    c5_l.fill.solid()
    c5_l.fill.fore_color.rgb = COLOR_CARD_BG
    c5_l.line.color.rgb = COLOR_CARD_BORDER
    tf5_l = c5_l.text_frame
    tf5_l.word_wrap = True
    tf5_l.margin_left = Inches(0.25)
    tf5_l.margin_top = Inches(0.2)
    tf5_l.margin_right = Inches(0.25)

    p = tf5_l.paragraphs[0]
    p.text = "• Potential impact on the target audience"
    p.font.size = FONT_SECTION_HEADER
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_FAMILY
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(10)

    auds = [
        ("Government & Defense Analysts (NTRO/MoD):", " Saves 4-6 hours daily per analyst by instantly synthesizing complex intelligence and security advisories into actionable briefs."),
        ("Senior Bureaucrats & C-Suite Executives:", " Enables immediate decision-making via 1-page executive memos and interactive Keynote presentation decks."),
        ("Public Relations & Media Officers:", " Instant generation of accurate, jargon-free press releases and infographics in multiple regional languages."),
        ("Field Personnel & Traveling Officials:", " Delivers 60-second audio podcast briefings directly to mobile devices.")
    ]
    for aud, text in auds:
        p = tf5_l.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {aud} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = FONT_FAMILY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(9.5)
        r2.font.name = FONT_FAMILY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(6)

    # Right Container: Benefits of the Solution
    c5_r = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.88) + col_5_w + col_5_gap, Inches(1.4), col_5_w, Inches(5.4))
    c5_r.fill.solid()
    c5_r.fill.fore_color.rgb = COLOR_CARD_BG
    c5_r.line.color.rgb = COLOR_CARD_BORDER
    tf5_r = c5_r.text_frame
    tf5_r.word_wrap = True
    tf5_r.margin_left = Inches(0.25)
    tf5_r.margin_top = Inches(0.2)
    tf5_r.margin_right = Inches(0.25)

    p = tf5_r.paragraphs[0]
    p.text = "• Benefits of the solution (social, economic, etc.)"
    p.font.size = FONT_SECTION_HEADER
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_FAMILY
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(10)

    bens = [
        ("💰 Economic Benefits:", " Reduces enterprise document processing costs by 75%; eliminates expensive third-party graphic and presentation outsourcing."),
        ("🤝 Social & Governance Impact:", " Bridges the communication gap between technical policymakers and citizens through regional Indic language summaries and visual posters."),
        ("🛡️ National Security Advantage:", " Maintains tactical information advantage during national security events through instant, multi-format threat dissemination."),
        ("🌱 Environmental Sustainability:", " Cuts down unnecessary paper report printing through interactive digital briefs and mobile voice summaries.")
    ]
    for lbl, text in bens:
        p = tf5_r.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{lbl} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = FONT_FAMILY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(9.5)
        r2.font.name = FONT_FAMILY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(6)

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES (Clean Citations & Verification)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_header_footer(slide6, 6, "RESEARCH AND REFERENCES")

    # Left Container: Research Papers
    c6_l = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.88), Inches(1.4), col_5_w, Inches(5.4))
    c6_l.fill.solid()
    c6_l.fill.fore_color.rgb = COLOR_CARD_BG
    c6_l.line.color.rgb = COLOR_CARD_BORDER
    tf6_l = c6_l.text_frame
    tf6_l.word_wrap = True
    tf6_l.margin_left = Inches(0.25)
    tf6_l.margin_top = Inches(0.2)
    tf6_l.margin_right = Inches(0.25)

    p = tf6_l.paragraphs[0]
    p.text = "• Details / Links of reference & research work"
    p.font.size = FONT_SECTION_HEADER
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_FAMILY
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(10)

    refs = [
        ("• Lewis et al. (2020):", "Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks (NeurIPS)."),
        ("• Touvron et al. (2023/2024):", "Llama 3 & Open Foundation Models for Sovereign Enterprise Reasoning."),
        ("• Vaswani et al. (2017):", "Attention Is All You Need — Transformer Architecture Fundamentals."),
        ("• AI4Bharat / Bhashini (2023):", "IndicBERT & Multilingual Natural Language Architectures for Indian Languages."),
        ("• ISO/IEC 27001 & NIST Standards:", "Cybersecurity & Data Protection Frameworks for Air-Gapped Intelligence.")
    ]
    for cit, title in refs:
        p = tf6_l.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{cit} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = FONT_FAMILY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = title
        r2.font.size = Pt(9.5)
        r2.font.name = FONT_FAMILY
        r2.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(6)

    # Right Container: Project & Official Validation Links
    c6_r = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.88) + col_5_w + col_5_gap, Inches(1.4), col_5_w, Inches(5.4))
    c6_r.fill.solid()
    c6_r.fill.fore_color.rgb = COLOR_CARD_BG
    c6_r.line.color.rgb = COLOR_CARD_BORDER
    tf6_r = c6_r.text_frame
    tf6_r.word_wrap = True
    tf6_r.margin_left = Inches(0.25)
    tf6_r.margin_top = Inches(0.2)
    tf6_r.margin_right = Inches(0.25)

    p = tf6_r.paragraphs[0]
    p.text = "• Official Portals & Project Verification Links"
    p.font.size = FONT_SECTION_HEADER
    p.font.bold = True
    p.font.underline = True
    p.font.name = FONT_FAMILY
    p.font.color.rgb = COLOR_PRIMARY_BLUE
    p.space_after = Pt(10)

    proj_links = [
        ("SIH 2026 Problem Statement 26154:", "https://www.sih.gov.in/sih2026PS"),
        ("NTRO Official Government Portal:", "https://ntro.gov.in"),
        ("Live Web Platform & Design System:", "https://hrl-brand-seo.vercel.app"),
        ("GitHub Project Repository:", "https://github.com/hrlpavan/sih-2026-media-problem-statements"),
        ("Winning Blueprint & Documentation:", "https://github.com/hrlpavan/sih-2026-media-problem-statements/blob/main/WINNING_STRATEGY_PS26154.md")
    ]
    for lbl, url in proj_links:
        p = tf6_r.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {lbl}\n  "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.name = FONT_FAMILY
        r1.font.color.rgb = COLOR_PRIMARY_BLUE
        r2 = p.add_run()
        r2.text = url
        r2.font.size = Pt(9.5)
        r2.font.name = FONT_FAMILY
        r2.font.color.rgb = COLOR_FOOTER_BLUE
        p.space_after = Pt(6)

    prs.save(output_path)
    print(f"Golden Ratio SIH 2026 deck created at: {output_path}")

if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "SIH2026_Idea_Presentation_PS26154.pptx"
    create_golden_ratio_sih_deck(out)
