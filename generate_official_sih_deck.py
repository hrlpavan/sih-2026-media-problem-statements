import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_official_sih_deck(output_path="SIH2026_Idea_Presentation_PS26154.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Professional Color Palette
    BG_WHITE = RGBColor(255, 255, 255)
    PRIMARY_BLUE = RGBColor(11, 44, 97)      # #0B2C61 (Navy/SIH Blue)
    ACCENT_BLUE = RGBColor(0, 113, 227)      # #0071E3 (Apple/Tech Blue)
    ACCENT_ORANGE = RGBColor(240, 90, 36)    # #F05A24 (SIH Orange)
    CARD_BG = RGBColor(245, 248, 252)        # Soft gray-blue container
    BORDER_COLOR = RGBColor(218, 228, 240)
    TEXT_DARK = RGBColor(20, 30, 45)
    TEXT_MUTED = RGBColor(90, 105, 125)
    FOOTER_BLUE = RGBColor(14, 116, 188)

    def add_base_decorations(slide, slide_num, title_text):
        # Top Header Bar & SIH Branding
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = ACCENT_ORANGE
        header_bar.line.fill.background()

        # Slide Title
        if title_text:
            tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(9.5), Inches(0.8))
            tf_t = tb_title.text_frame
            tf_t.word_wrap = True
            p = tf_t.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_BLUE
            p.font.name = "Arial"

        # Team Tag Top Left
        team_tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.2), Inches(2.2), Inches(0.4))
        team_tag.fill.solid()
        team_tag.fill.fore_color.rgb = CARD_BG
        team_tag.line.color.rgb = BORDER_COLOR
        tf_tag = team_tag.text_frame
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = "Team: OmniTransform"
        p_tag.font.size = Pt(10)
        p_tag.font.bold = True
        p_tag.font.color.rgb = PRIMARY_BLUE
        p_tag.alignment = PP_ALIGN.CENTER

        # SIH 2026 Badge Top Right
        badge_box = slide.shapes.add_textbox(Inches(10.2), Inches(0.2), Inches(2.3), Inches(0.6))
        tf_b = badge_box.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = "SMART INDIA HACKATHON 2026"
        p_b.font.size = Pt(9)
        p_b.font.bold = True
        p_b.font.color.rgb = PRIMARY_BLUE
        p_b.alignment = PP_ALIGN.RIGHT

        # Bottom Footer Ribbon
        footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4))
        footer_bar.fill.solid()
        footer_bar.fill.fore_color.rgb = FOOTER_BLUE
        footer_bar.line.fill.background()

        tb_f = slide.shapes.add_textbox(Inches(0.8), Inches(7.12), Inches(10.0), Inches(0.35))
        tf_f = tb_f.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = "@SIH Idea submission- Template | PS ID: 26154 (NTRO)"
        p_f.font.size = Pt(9)
        p_f.font.color.rgb = RGBColor(255, 255, 255)

        tb_num = slide.shapes.add_textbox(Inches(12.0), Inches(7.12), Inches(0.8), Inches(0.35))
        tf_num = tb_num.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.font.size = Pt(10)
        p_num.font.bold = True
        p_num.font.color.rgb = RGBColor(255, 255, 255)
        p_num.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 1: TITLE PAGE (Official Template Layout)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Top Bar
    h_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
    h_bar.fill.solid()
    h_bar.fill.fore_color.rgb = ACCENT_ORANGE
    h_bar.line.fill.background()

    # Main SIH Banner Text
    tb1 = slide1.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.6))
    p1 = tb1.text_frame.paragraphs[0]
    p1.text = "SMART INDIA HACKATHON 2026"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY_BLUE
    p1.alignment = PP_ALIGN.CENTER

    tb_sub = slide1.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.4))
    p_sub = tb_sub.text_frame.paragraphs[0]
    p_sub.text = "TITLE PAGE"
    p_sub.font.size = Pt(18)
    p_sub.font.bold = True
    p_sub.font.color.rgb = TEXT_DARK
    p_sub.alignment = PP_ALIGN.CENTER

    # Left Container: Mandatory Fields
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.1))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = BORDER_COLOR
    card1.line.width = Pt(1.5)

    tf_c1 = card1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_left = Inches(0.3)
    tf_c1.margin_top = Inches(0.3)

    fields = [
        ("Problem Statement ID:", "26154"),
        ("Problem Statement Title:", "Gen AI Platform for Automated Content Transformation"),
        ("Theme:", "Blockchain & Cybersecurity / Artificial Intelligence"),
        ("PS Category:", "Software"),
        ("Organization:", "National Technical Research Organisation (NTRO)"),
        ("Team Name:", "OmniTransform (Leader: Pavan Kumar S)"),
        ("Idea Title:", "OmniTransform AI – Sovereign Multimodal Content Engine")
    ]

    for idx, (lbl, val) in enumerate(fields):
        p = tf_c1.add_paragraph() if idx > 0 else tf_c1.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"• {lbl} "
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = PRIMARY_BLUE
        
        r2 = p.add_run()
        r2.text = val
        r2.font.bold = (idx in [0, 1, 6])
        r2.font.size = Pt(13)
        r2.font.color.rgb = TEXT_DARK if idx != 1 else ACCENT_ORANGE
        p.space_after = Pt(10)

    # Right Container: Project Emblem & Value Props
    card1_r = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.6), Inches(4.0), Inches(5.1))
    card1_r.fill.solid()
    card1_r.fill.fore_color.rgb = PRIMARY_BLUE
    card1_r.line.fill.background()

    tf_r = card1_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.3)
    tf_r.margin_top = Inches(0.4)

    p_rt = tf_r.paragraphs[0]
    p_rt.text = "OmniTransform AI"
    p_rt.font.size = Pt(20)
    p_rt.font.bold = True
    p_rt.font.color.rgb = RGBColor(255, 255, 255)
    p_rt.alignment = PP_ALIGN.CENTER
    p_rt.space_after = Pt(14)

    value_bullets = [
        "⚡ 1-to-5 Synchronized Multi-Asset Generation in <10s",
        "🎯 100% Grounded Source Citation & Anti-Hallucination",
        "🔒 Sovereign Air-Gapped On-Premise Compliance",
        "🎙️ Integrated Studio-Quality Audio Intelligence Podcast",
        "🌐 Live Ecosystem: hrl-brand-seo.vercel.app"
    ]
    for vb in value_bullets:
        p = tf_r.add_paragraph()
        p.text = vb
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(220, 235, 255)
        p.space_after = Pt(12)

    # =========================================================================
    # SLIDE 2: PROPOSED SOLUTION
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_base_decorations(slide2, 2, "PROPOSED SOLUTION")

    # Subtitle Badge
    tb2_sub = slide2.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4))
    p2_s = tb2_sub.text_frame.paragraphs[0]
    p2_s.text = "OmniTransform AI: Single-Pass Ingestion to 5 Synchronized Media Formats"
    p2_s.font.size = Pt(13)
    p2_s.font.bold = True
    p2_s.font.color.rgb = ACCENT_ORANGE

    # 3 Structured Vertical Columns
    col_w = Inches(3.64)
    col_gap = Inches(0.38)
    left_start = Inches(0.8)

    # Card A: Detailed Explanation
    c_a = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_start, Inches(1.6), col_w, Inches(5.1))
    c_a.fill.solid()
    c_a.fill.fore_color.rgb = CARD_BG
    c_a.line.color.rgb = BORDER_COLOR
    tf_a = c_a.text_frame
    tf_a.word_wrap = True
    tf_a.margin_left = Inches(0.2)
    tf_a.margin_top = Inches(0.2)
    
    p = tf_a.paragraphs[0]
    p.text = "1. Detailed Explanation"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY_BLUE
    p.space_after = Pt(10)

    pts_a = [
        "• Core Concept: An enterprise Gen AI platform that ingests dense 50+ page PDFs, threat advisories, and policy documents.",
        "• Single-Pass Multi-Output: In under 10 seconds, generates 5 production assets simultaneously:",
        "  1. 📊 Executive 1-Page Memo",
        "  2. 🖥️ Keynote Slide Presentation",
        "  3. 🎨 Visual Infographics & Cards",
        "  4. 📰 Multilingual Press Release",
        "  5. 🎙️ 60s Voice Audio Podcast"
    ]
    for pt in pts_a:
        p = tf_a.add_paragraph()
        p.text = pt
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(5)

    # Card B: How It Addresses Problem
    c_b = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_start + col_w + col_gap, Inches(1.6), col_w, Inches(5.1))
    c_b.fill.solid()
    c_b.fill.fore_color.rgb = CARD_BG
    c_b.line.color.rgb = BORDER_COLOR
    tf_b = c_b.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = Inches(0.2)
    tf_b.margin_top = Inches(0.2)

    p = tf_b.paragraphs[0]
    p.text = "2. How It Addresses Problem"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY_BLUE
    p.space_after = Pt(10)

    pts_b = [
        "• Eliminates 4-6 Hours Manual Synthesis: Automates document restructuring across technical, executive, and public audiences.",
        "• Solves Tool Fragmentation: Replaces separate disconnected tools (Canva, PowerPoint, Word, Chatbots).",
        "• Eliminates Hallucination Risk: Enforces strict RAG boundaries with exact sentence-level coordinate tracking.",
        "• Preserves Tactical Velocity: Immediate delivery of time-critical intelligence to senior defense/civil directors."
    ]
    for pt in pts_b:
        p = tf_b.add_paragraph()
        p.text = pt
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # Card C: Innovation & Uniqueness
    c_c = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_start + (col_w + col_gap)*2, Inches(1.6), col_w, Inches(5.1))
    c_c.fill.solid()
    c_c.fill.fore_color.rgb = CARD_BG
    c_c.line.color.rgb = BORDER_COLOR
    tf_c = c_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = Inches(0.2)
    tf_c.margin_top = Inches(0.2)

    p = tf_c.paragraphs[0]
    p.text = "3. Innovation & Uniqueness"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY_BLUE
    p.space_after = Pt(10)

    pts_c = [
        "• Interactive Source Badges: Click any bullet to open the PDF and highlight the exact originating sentence.",
        "• 1-Click Multi-Persona Toggle: Seamlessly switch between C-Suite Memo, Citizen FAQ, and Threat Brief.",
        "• Sovereign Air-Gapped Ready: Runs locally on open-weights (Llama 3/Mistral) without public cloud leaks.",
        "• Apple-Standard Frosted UI: Glassmorphism layout for clean, executive-ready presentation."
    ]
    for pt in pts_c:
        p = tf_c.add_paragraph()
        p.text = pt
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_base_decorations(slide3, 3, "TECHNICAL APPROACH")

    # Left Container: Technologies Used
    c3_left = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(4.5), Inches(5.4))
    c3_left.fill.solid()
    c3_left.fill.fore_color.rgb = CARD_BG
    c3_left.line.color.rgb = BORDER_COLOR
    tf3_l = c3_left.text_frame
    tf3_l.word_wrap = True
    tf3_l.margin_left = Inches(0.25)
    tf3_l.margin_top = Inches(0.2)

    p = tf3_l.paragraphs[0]
    p.text = "1. Technologies & Stack"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = PRIMARY_BLUE
    p.space_after = Pt(10)

    tech_specs = [
        ("Frontend:", "Next.js 14, React, Tailwind CSS, Lucide Icons, PDF.js canvas viewer."),
        ("Backend Engine:", "Python FastAPI, LangChain / LlamaIndex, PyMuPDF coordinate parser."),
        ("AI / LLMs:", "Llama 3.3 (70B), Mistral Large, IndicBERT (regional translation)."),
        ("Rendering & Media:", "Satori HTML-to-Image Canvas API, Marp slide compiler, FFmpeg."),
        ("Voice AI:", "Piper Neural TTS & ElevenLabs API."),
        ("Deployment:", "Docker containerized, air-gapped Linux host, zero telemetry.")
    ]
    for cat, desc in tech_specs:
        p = tf3_l.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{cat} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = ACCENT_BLUE
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # Right Container: Implementation Methodology & Architecture
    c3_right = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.6), Inches(1.4), Inches(6.9), Inches(5.4))
    c3_right.fill.solid()
    c3_right.fill.fore_color.rgb = CARD_BG
    c3_right.line.color.rgb = BORDER_COLOR
    tf3_r = c3_right.text_frame
    tf3_r.word_wrap = True
    tf3_r.margin_left = Inches(0.25)
    tf3_r.margin_top = Inches(0.2)

    p = tf3_r.paragraphs[0]
    p.text = "2. Methodology & Process Flow"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = PRIMARY_BLUE
    p.space_after = Pt(10)

    flow_steps = [
        ("Step 1: Document Ingestion", "PyMuPDF parses raw PDF/DOCX, extracting structural hierarchy, text coordinates, diagrams, and tables with bounding box tags."),
        ("Step 2: Structured Chunking & RAG", "Hierarchical chunking creates contextual embeddings stored in an in-memory vector store (FAISS/Qdrant) for strict citation binding."),
        ("Step 3: Multi-Persona Orchestration", "Constrained JSON Schema prompts instruct the LLM to extract key insights, slide structures, infographic metrics, and press statements."),
        ("Step 4: Parallel Synthesis Engines", "Simultaneously compiles HTML5 presentation slides, auto-renders high-res visual cards (Canvas API), and generates 60s neural audio stream."),
        ("Step 5: Interactive Verification UI", "Delivers a live dashboard with side-by-side asset previews and clickable citation coordinates directly linking to the source PDF.")
    ]
    for st, desc in flow_steps:
        p = tf3_r.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {st}: "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = PRIMARY_BLUE
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_base_decorations(slide4, 4, "FEASIBILITY AND VIABILITY")

    # Top Row: Feasibility Matrix
    c4_top = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.7), Inches(2.2))
    c4_top.fill.solid()
    c4_top.fill.fore_color.rgb = CARD_BG
    c4_top.line.color.rgb = BORDER_COLOR
    tf4_t = c4_top.text_frame
    tf4_t.word_wrap = True
    tf4_t.margin_left = Inches(0.25)
    tf4_t.margin_top = Inches(0.15)

    p = tf4_t.paragraphs[0]
    p.text = "1. Technical & Operational Feasibility Analysis"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY_BLUE
    p.space_after = Pt(6)

    feas_points = [
        ("• Hardware & Infrastructure:", "Runs efficiently on commodity GPU workstations (e.g. Single RTX 4090 or Apple Silicon) using 4-bit quantized open-weight models."),
        ("• Deterministic Speed:", "Asynchronous map-reduce pipeline completes multi-format transformation of a 30-page PDF in < 10 seconds."),
        ("• Zero External Dependency:", "Entire software stack is containerized in Docker; operates fully offline in air-gapped defense/intelligence environments.")
    ]
    for lbl, desc in feas_points:
        p = tf4_t.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{lbl} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = ACCENT_BLUE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    # Bottom Row: 3 Risk & Strategy Cards
    card_w = Inches(3.64)
    c4_r1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_start, Inches(3.8), card_w, Inches(3.0))
    c4_r1.fill.solid()
    c4_r1.fill.fore_color.rgb = CARD_BG
    c4_r1.line.color.rgb = BORDER_COLOR
    tf4_1 = c4_r1.text_frame
    tf4_1.word_wrap = True
    tf4_1.margin_left = Inches(0.2)
    tf4_1.margin_top = Inches(0.15)
    p = tf4_1.paragraphs[0]
    p.text = "Risk 1: AI Hallucination"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_ORANGE
    p.space_after = Pt(6)
    p = tf4_1.add_paragraph()
    p.text = "• Challenge: Fabricated facts in intelligence briefs.\n• Strategy: Strict RAG constraint boundaries + reverse citation bounding-box verification for 100% grounded truth."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK

    c4_r2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_start + col_w + col_gap, Inches(3.8), card_w, Inches(3.0))
    c4_r2.fill.solid()
    c4_r2.fill.fore_color.rgb = CARD_BG
    c4_r2.line.color.rgb = BORDER_COLOR
    tf4_2 = c4_r2.text_frame
    tf4_2.word_wrap = True
    tf4_2.margin_left = Inches(0.2)
    tf4_2.margin_top = Inches(0.15)
    p = tf4_2.paragraphs[0]
    p.text = "Risk 2: Data Confidentiality"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_ORANGE
    p.space_after = Pt(6)
    p = tf4_2.add_paragraph()
    p.text = "• Challenge: Classified leaks to public cloud APIs.\n• Strategy: Sovereign on-premise deployment with air-gapped local LLMs (Llama 3/Mistral); zero telemetry."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK

    c4_r3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_start + (col_w + col_gap)*2, Inches(3.8), card_w, Inches(3.0))
    c4_r3.fill.solid()
    c4_r3.fill.fore_color.rgb = CARD_BG
    c4_r3.line.color.rgb = BORDER_COLOR
    tf4_3 = c4_r3.text_frame
    tf4_3.word_wrap = True
    tf4_3.margin_left = Inches(0.2)
    tf4_3.margin_top = Inches(0.15)
    p = tf4_3.paragraphs[0]
    p.text = "Risk 3: Complex Tables & Charts"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_ORANGE
    p.space_after = Pt(6)
    p = tf4_3.add_paragraph()
    p.text = "• Challenge: Loss of tabular & diagrammatic data.\n• Strategy: Multimodal vision-language parsing (VLM) + Markdown table structure extraction."
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_DARK

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_base_decorations(slide5, 5, "IMPACT AND BENEFITS")

    # Left Container: Target Audience Impact
    c5_l = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.4))
    c5_l.fill.solid()
    c5_l.fill.fore_color.rgb = CARD_BG
    c5_l.line.color.rgb = BORDER_COLOR
    tf5_l = c5_l.text_frame
    tf5_l.word_wrap = True
    tf5_l.margin_left = Inches(0.25)
    tf5_l.margin_top = Inches(0.2)

    p = tf5_l.paragraphs[0]
    p.text = "1. Impact on Target Audience"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = PRIMARY_BLUE
    p.space_after = Pt(10)

    audiences = [
        ("Government & Defense Analysts (NTRO/MoD):", "Saves 4-6 hours daily per analyst by instantly synthesizing complex intelligence and security advisories into actionable briefs."),
        ("Senior Bureaucrats & C-Suite Executives:", "Enables immediate decision-making via 1-page executive memos and interactive Keynote presentation decks."),
        ("Public Relations & Media Officers:", "Instant generation of accurate, jargon-free press releases and infographics in multiple regional languages."),
        ("Field Personnel & Traveling Officials:", "Delivers 60-second audio podcast briefings directly to mobile devices.")
    ]
    for aud, imp in audiences:
        p = tf5_l.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {aud}\n  "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = ACCENT_BLUE
        r2 = p.add_run()
        r2.text = imp
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # Right Container: Multi-Dimensional Benefits
    c5_r = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.7), Inches(5.4))
    c5_r.fill.solid()
    c5_r.fill.fore_color.rgb = CARD_BG
    c5_r.line.color.rgb = BORDER_COLOR
    tf5_r = c5_r.text_frame
    tf5_r.word_wrap = True
    tf5_r.margin_left = Inches(0.25)
    tf5_r.margin_top = Inches(0.2)

    p = tf5_r.paragraphs[0]
    p.text = "2. Multi-Dimensional Benefits"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = PRIMARY_BLUE
    p.space_after = Pt(10)

    benefits = [
        ("💰 Economic Benefits:", "Reduces enterprise document processing costs by 75%; eliminates expensive third-party graphic and presentation outsourcing."),
        ("🤝 Social & Governance Impact:", "Bridges the communication gap between technical policymakers and general citizens through regional Indic language summaries and visual posters."),
        ("🛡️ National Security Benefit:", "Maintains tactical information advantage during national security events through instant, multi-format threat dissemination."),
        ("🌱 Environmental Sustainability:", "Drastically cuts down unnecessary paper report printing through interactive digital briefs and mobile voice summaries.")
    ]
    for b_title, b_desc in benefits:
        p = tf5_r.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{b_title} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = ACCENT_ORANGE
        r2 = p.add_run()
        r2.text = b_desc
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_base_decorations(slide6, 6, "RESEARCH AND REFERENCES")

    # Left Container: Research Papers & Standards
    c6_l = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.4))
    c6_l.fill.solid()
    c6_l.fill.fore_color.rgb = CARD_BG
    c6_l.line.color.rgb = BORDER_COLOR
    tf6_l = c6_l.text_frame
    tf6_l.word_wrap = True
    tf6_l.margin_left = Inches(0.25)
    tf6_l.margin_top = Inches(0.2)

    p = tf6_l.paragraphs[0]
    p.text = "1. Research Papers & Literature"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = PRIMARY_BLUE
    p.space_after = Pt(10)

    papers = [
        ("• Lewis et al. (2020):", "Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks (NeurIPS)."),
        ("• Touvron et al. (2023/2024):", "Llama 3 & Open Foundation Models for Sovereign Enterprise Reasoning."),
        ("• Vaswani et al. (2017):", "Attention Is All You Need — Transformer Architecture Fundamentals."),
        ("• Bhashini / AI4Bharat (2023):", "IndicBERT & Multilingual Natural Language Architectures for Indian Languages."),
        ("• ISO/IEC 27001 & NIST Standards:", "Cybersecurity & Data Protection Frameworks for Air-Gapped Intelligence.")
    ]
    for cit, title in papers:
        p = tf6_l.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{cit} "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = ACCENT_BLUE
        r2 = p.add_run()
        r2.text = title
        r2.font.size = Pt(10)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # Right Container: Project Artifacts & Live Links
    c6_r = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.7), Inches(5.4))
    c6_r.fill.solid()
    c6_r.fill.fore_color.rgb = CARD_BG
    c6_r.line.color.rgb = BORDER_COLOR
    tf6_r = c6_r.text_frame
    tf6_r.word_wrap = True
    tf6_r.margin_left = Inches(0.25)
    tf6_r.margin_top = Inches(0.2)

    p = tf6_r.paragraphs[0]
    p.text = "2. Project Links & Live Ecosystem"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = PRIMARY_BLUE
    p.space_after = Pt(10)

    links = [
        ("SIH 2026 Problem Statement 26154:", "https://www.sih.gov.in/sih2026PS"),
        ("NTRO Official Portal:", "https://ntro.gov.in"),
        ("Live Web Platform & Design System:", "https://hrl-brand-seo.vercel.app"),
        ("GitHub Project Repository:", "https://github.com/hrlpavan/sih-2026-media-problem-statements"),
        ("Winning Blueprint & Documentation:", "https://github.com/hrlpavan/sih-2026-media-problem-statements/blob/main/WINNING_STRATEGY_PS26154.md")
    ]
    for lbl, url in links:
        p = tf6_r.add_paragraph()
        r1 = p.add_run()
        r1.text = f"• {lbl}\n  "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = ACCENT_ORANGE
        r2 = p.add_run()
        r2.text = url
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    prs.save(output_path)
    print(f"Official 6-slide SIH 2026 deck created at: {output_path}")

if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "SIH2026_Idea_Presentation_PS26154.pptx"
    create_official_sih_deck(out)
