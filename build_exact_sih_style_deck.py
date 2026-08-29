import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_exact_sih_deck(output_path="SIH2026_Idea_Presentation_PS26154.pptx"):
    prs = Presentation()
    # Official 16:9 Widescreen (13.333 in x 7.5 in)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Colors from the exact reference sample
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_NAVY = RGBColor(27, 54, 93)          # Primary Title Navy Blue (#1B365D)
    COLOR_GREEN_HDR = RGBColor(106, 168, 79)    # Reference Green Header (#6AA84F)
    COLOR_TEAL_HDR = RGBColor(60, 120, 140)     # Reference Teal/Blue Header (#3C788C)
    COLOR_ORANGE_HDR = RGBColor(230, 145, 56)   # Reference Orange/Gold Header (#E69138)
    COLOR_FOOTER_BLUE = RGBColor(0, 119, 200)   # Official SIH Footer Ribbon (#0077C8)
    COLOR_TEXT = RGBColor(20, 20, 20)           # Clean dark text
    COLOR_CARD_BG = RGBColor(245, 248, 250)     # Soft subtle card fill
    COLOR_CARD_BORDER = RGBColor(210, 220, 230) # Soft border

    FONT_TITLE = "Times New Roman"
    FONT_BODY = "Arial"
    LOGO_PATH = "sih_official_logo.png"
    FLOWCHART_PATH = "omnitransform_pipeline_flowchart.png"
    
    # DEDICATED MASTER REPOSITORY & RESOURCES URL
    DEDICATED_REPO_URL = "https://github.com/hrlpavan/omnitransform-ai-resources"

    def add_header_footer(slide, slide_num, center_title, team_name="OmniTransform"):
        # Top Left Team Oval
        team_oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.4), Inches(0.2), Inches(1.8), Inches(0.95))
        team_oval.fill.solid()
        team_oval.fill.fore_color.rgb = COLOR_WHITE
        team_oval.line.color.rgb = COLOR_NAVY
        team_oval.line.width = Pt(1.5)
        tf_team = team_oval.text_frame
        tf_team.word_wrap = True
        tf_team.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf_team.paragraphs[0]
        p1.text = team_name
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.name = FONT_BODY
        p1.font.color.rgb = COLOR_NAVY
        p1.alignment = PP_ALIGN.CENTER

        # Top Center Title
        if center_title:
            tb_title = slide.shapes.add_textbox(Inches(2.3), Inches(0.25), Inches(7.4), Inches(0.85))
            tf_title = tb_title.text_frame
            tf_title.word_wrap = True
            p_t = tf_title.paragraphs[0]
            p_t.text = center_title
            p_t.font.size = Pt(22 if len(center_title) > 25 else 26)
            p_t.font.bold = True
            p_t.font.name = FONT_TITLE
            p_t.font.color.rgb = COLOR_NAVY
            p_t.alignment = PP_ALIGN.CENTER

        # Top Right SIH Logo
        if os.path.exists(LOGO_PATH):
            slide.shapes.add_picture(LOGO_PATH, Inches(10.0), Inches(0.18), width=Inches(2.8))

        # Bottom Footer Bar
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), Inches(13.333), Inches(0.45))
        footer.fill.solid()
        footer.fill.fore_color.rgb = COLOR_FOOTER_BLUE
        footer.line.fill.background()

        tb_f = slide.shapes.add_textbox(Inches(0.8), Inches(7.08), Inches(10.0), Inches(0.38))
        p_f = tb_f.text_frame.paragraphs[0]
        p_f.text = "@SIH Idea submission- Template"
        p_f.font.size = Pt(10)
        p_f.font.name = FONT_BODY
        p_f.font.color.rgb = COLOR_WHITE

        tb_num = slide.shapes.add_textbox(Inches(12.0), Inches(7.08), Inches(0.8), Inches(0.38))
        p_num = tb_num.text_frame.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.font.size = Pt(11)
        p_num.font.bold = True
        p_num.font.name = FONT_BODY
        p_num.font.color.rgb = COLOR_WHITE
        p_num.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 1: TITLE PAGE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)

    tb1_top = slide1.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7))
    p_t1 = tb1_top.text_frame.paragraphs[0]
    p_t1.text = "SMART INDIA HACKATHON 2026"
    p_t1.font.size = Pt(32)
    p_t1.font.bold = True
    p_t1.font.name = FONT_TITLE
    p_t1.font.color.rgb = COLOR_NAVY
    p_t1.alignment = PP_ALIGN.LEFT

    tb1_fields = slide1.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.4))
    tf1 = tb1_fields.text_frame
    tf1.word_wrap = True

    fields = [
        ("• Problem Statement ID – ", "26154"),
        ("• Problem Statement Title - ", "Gen AI Platform for Automated Content Transformation"),
        ("• Theme - ", "Miscellaneous / AI"),
        ("• PS Category - ", "Software"),
        ("• Team ID - ", "104580"),
        ("• Team Name - ", "OmniTransform"),
        ("• Project Resources - ", DEDICATED_REPO_URL)
    ]

    for idx, (lbl, val) in enumerate(fields):
        p = tf1.add_paragraph() if idx > 0 else tf1.paragraphs[0]
        r1 = p.add_run()
        r1.text = lbl
        r1.font.bold = True
        r1.font.size = Pt(16 if idx == 6 else 17.5)
        r1.font.name = FONT_TITLE
        r1.font.color.rgb = COLOR_NAVY if idx == 6 else COLOR_TEXT

        r2 = p.add_run()
        r2.text = val
        r2.font.bold = (idx in [0, 1, 5, 6])
        r2.font.size = Pt(14 if idx == 6 else 17.5)
        r2.font.underline = (idx == 6)
        r2.font.name = FONT_BODY
        r2.font.color.rgb = COLOR_FOOTER_BLUE if idx == 6 else COLOR_TEXT
        p.space_after = Pt(10)

    if os.path.exists(LOGO_PATH):
        slide1.shapes.add_picture(LOGO_PATH, Inches(9.2), Inches(1.8), width=Inches(3.6))

    # =========================================================================
    # SLIDE 2: PROPOSED SOLUTION
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide2, 2, "OMNITRANSFORM (ALL-IN-ONE) AI CONTENT PLATFORM")

    col_w = Inches(3.8)
    col_gap = Inches(0.3)
    c_top = Inches(1.3)
    c_h = Inches(5.5)

    def create_column_card(left_in, header_title, header_color, bullet_items):
        hdr = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_in, c_top, col_w, Inches(0.45))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = header_color
        hdr.line.fill.background()
        tf_h = hdr.text_frame
        tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf_h.paragraphs[0]
        p.text = header_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.name = FONT_BODY
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

        card = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_in, c_top + Inches(0.45), col_w, c_h - Inches(0.45))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_CARD_BORDER
        card.line.width = Pt(1)
        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.18)
        tf_c.margin_top = Inches(0.18)
        tf_c.margin_right = Inches(0.18)

        for idx, (title, desc) in enumerate(bullet_items):
            p = tf_c.add_paragraph() if idx > 0 else tf_c.paragraphs[0]
            r1 = p.add_run()
            r1.text = f"● {title}: "
            r1.font.bold = True
            r1.font.size = Pt(11)
            r1.font.name = FONT_BODY
            r1.font.color.rgb = COLOR_TEXT

            r2 = p.add_run()
            r2.text = desc
            r2.font.size = Pt(10.5)
            r2.font.name = FONT_BODY
            r2.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(7)

    col1_items = [
        ("Simple Web Platform", "A clean and fast web application built for government officers, analysts, and students."),
        ("Upload Any Document", "Takes long PDFs, research whitepapers, and government reports up to 100+ pages."),
        ("5 Output Formats in Seconds", "Automatically creates 1-page memo, presentation slides, infographics, press release, and 60s voice audio."),
        ("Saves Hours of Manual Work", "Officers save 4 to 6 hours daily spent on manually typing summaries and designing PPTs."),
        ("Accurate Page References", "Every bullet point links directly to the exact page and line in the original PDF file."),
        ("Works 100% Offline", "Can be installed on local office servers without sending sensitive data to the internet.")
    ]
    create_column_card(Inches(0.6), "Proposed Solution Overview", COLOR_GREEN_HDR, col1_items)

    col2_items = [
        ("Smart PDF Parser", "Extracts clean text, headings, and complex tables using PyMuPDF and pdfplumber."),
        ("Fast AI Summarization", "Uses open-source Llama 3 and Mistral models to extract key decisions and highlights."),
        ("Slide Deck Creator", "Converts raw notes into clean, formatted PowerPoint slides ready for team meetings."),
        ("Multilingual Press Release", "Translates official summaries into Hindi, English, and regional Indian languages."),
        ("60-Second Audio Podcast", "Generates clear voice audio briefings using local Piper neural text-to-speech."),
        ("Safe & Private Database", "Stores files in a local secure database with zero third-party cloud tracking.")
    ]
    create_column_card(Inches(0.6) + col_w + col_gap, "Key Feature & Technology", COLOR_TEAL_HDR, col2_items)

    col3_items = [
        ("All-in-One Solution", "Replaces 5 separate tools into one simple dashboard (Summary + PPT + Audio + Image + Translate)."),
        ("Zero Hallucinations", "The AI is locked to only use facts from the uploaded document, preventing fake numbers."),
        ("Click-to-Verify Citations", "Clicking any summary sentence instantly opens the PDF with the exact sentence highlighted."),
        ("Regional Language Reach", "Helps common citizens read complex government schemes in their own local language."),
        ("Fast Decision Making", "Leaders and senior officers can review urgent threat alerts in under 1 minute."),
        ("Low Cost & Open Source", "Built on free open-source software with no expensive monthly API bills.")
    ]
    create_column_card(Inches(0.6) + (col_w + col_gap)*2, "Innovation & Social Impact", COLOR_GREEN_HDR, col3_items)

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide3, 3, "TECHNICAL APPROACH")

    if os.path.exists(FLOWCHART_PATH):
        slide3.shapes.add_picture(FLOWCHART_PATH, Inches(0.6), Inches(1.3), width=Inches(6.6))

    tb3_r = slide3.shapes.add_textbox(Inches(7.4), Inches(1.3), Inches(5.3), Inches(4.7))
    tf3 = tb3_r.text_frame
    tf3.word_wrap = True

    p = tf3.paragraphs[0]
    p.text = "• Technologies that are used are Programming languages like"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.name = FONT_BODY
    p.font.color.rgb = COLOR_TEXT
    p.space_after = Pt(4)

    tech_list = [
        "JavaScript with React / Next.js (Web app development),",
        "Python FastAPI (Backend server & API),",
        "PyMuPDF & pdfplumber (PDF text & table extraction),",
        "Llama 3 / Mistral (AI summarization & text processing),",
        "IndicBERT (Multiple Indian Language Translation),",
        "SQLite / Vector DB (Database)."
    ]
    for t in tech_list:
        pt = tf3.add_paragraph()
        pt.text = f"  {t}"
        pt.font.size = Pt(11)
        pt.font.name = FONT_BODY
        pt.font.color.rgb = COLOR_TEXT
        pt.space_after = Pt(2)

    p_m = tf3.add_paragraph()
    p_m.text = "\n• Methodology and process for implementation :"
    p_m.font.size = Pt(13)
    p_m.font.bold = True
    p_m.font.name = FONT_BODY
    p_m.font.color.rgb = COLOR_TEXT
    p_m.space_after = Pt(4)

    steps = [
        "1) User uploads PDF / Document to web dashboard",
        "2) Parser extracts text, headings, and data tables",
        "3) AI creates 1-page summary & formatted slide deck",
        "4) Audio engine synthesizes 60s voice podcast briefing",
        "5) User previews live and downloads all 5 files"
    ]
    for s in steps:
        ps = tf3.add_paragraph()
        ps.text = f"  {s}"
        ps.font.size = Pt(11)
        ps.font.name = FONT_BODY
        ps.font.color.rgb = COLOR_TEXT
        ps.space_after = Pt(2)

    # Clickable Live Prototype Demo Box with Dedicated Repo URL
    tb3_link = slide3.shapes.add_textbox(Inches(7.4), Inches(6.1), Inches(5.3), Inches(0.85))
    tf_l = tb3_link.text_frame
    p_link = tf_l.paragraphs[0]
    p_link.text = "APP PROTOTYPE & RESOURCE REPOSITORY LINK"
    p_link.font.size = Pt(13.5)
    p_link.font.bold = True
    p_link.font.name = FONT_BODY
    p_link.font.color.rgb = COLOR_FOOTER_BLUE
    
    p_url = tf_l.add_paragraph()
    p_url.text = DEDICATED_REPO_URL
    p_url.font.size = Pt(12)
    p_url.font.bold = True
    p_url.font.underline = True
    p_url.font.name = FONT_BODY
    p_url.font.color.rgb = COLOR_FOOTER_BLUE

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide4, 4, "FEASIBILITY AND VIABILITY")

    tb4_l = slide4.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(7.5), Inches(5.5))
    tf4 = tb4_l.text_frame
    tf4.word_wrap = True

    p = tf4.paragraphs[0]
    p.text = "ANALYSIS OF THE FEASIBILITY OF THE IDEA"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.name = FONT_BODY
    p.font.color.rgb = COLOR_TEXT
    p.space_after = Pt(4)

    feas_points = [
        ("Technical & Practical Feasibility", "Built using standard open-source tools like Python, React, and local AI models. Works smoothly on a standard office laptop or PC."),
        ("Easy Offline Deployment", "The complete software runs inside Docker containers on local office computers without needing an active internet connection.")
    ]
    for title, desc in feas_points:
        pf = tf4.add_paragraph()
        r1 = pf.add_run()
        r1.text = f"● {title}: "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = COLOR_TEXT
        r2 = pf.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = COLOR_TEXT
        pf.space_after = Pt(5)

    pc = tf4.add_paragraph()
    pc.text = "\nPOTENTIAL CHALLENGES"
    pc.font.size = Pt(14)
    pc.font.bold = True
    pc.font.name = FONT_BODY
    pc.font.color.rgb = COLOR_TEXT
    pc.space_after = Pt(4)

    chall_points = [
        ("Complex Tables & Charts", "Dense multi-column tables in scanned PDFs can be hard to parse accurately."),
        ("Processing Speed on Long Files", "Converting 100+ page documents requires efficient processing to stay under 10 seconds.")
    ]
    for title, desc in chall_points:
        p_c = tf4.add_paragraph()
        r1 = p_c.add_run()
        r1.text = f"● {title}: "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p_c.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = COLOR_TEXT
        p_c.space_after = Pt(5)

    ps = tf4.add_paragraph()
    ps.text = "\nSTRATEGIES TO OVERCOME CHALLENGES"
    ps.font.size = Pt(14)
    ps.font.bold = True
    ps.font.name = FONT_BODY
    ps.font.color.rgb = COLOR_TEXT
    ps.space_after = Pt(4)

    strat_points = [
        ("Smart Table Parser", "Use layout-aware coordinate parsers (pdfplumber) to keep table rows and numbers intact."),
        ("Parallel Page Processing", "Split chapters into parallel worker tasks so long reports are summarized in seconds.")
    ]
    for title, desc in strat_points:
        p_s = tf4.add_paragraph()
        r1 = p_s.add_run()
        r1.text = f"● {title}: "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p_s.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = COLOR_TEXT
        p_s.space_after = Pt(5)

    c4_r = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.4), Inches(1.3), Inches(4.3), Inches(5.4))
    c4_r.fill.solid()
    c4_r.fill.fore_color.rgb = COLOR_CARD_BG
    c4_r.line.color.rgb = COLOR_CARD_BORDER
    tf4_r = c4_r.text_frame
    tf4_r.word_wrap = True
    tf4_r.margin_left = Inches(0.25)
    tf4_r.margin_top = Inches(0.3)
    tf4_r.margin_right = Inches(0.25)

    p_rt = tf4_r.paragraphs[0]
    p_rt.text = "KEY FEASIBILITY HIGHLIGHTS"
    p_rt.font.size = Pt(13)
    p_rt.font.bold = True
    p_rt.font.name = FONT_BODY
    p_rt.font.color.rgb = COLOR_NAVY
    p_rt.alignment = PP_ALIGN.CENTER
    p_rt.space_after = Pt(14)

    r_boxes = [
        ("[+] Low Cost", "Runs on free open-source models with no recurring software license fees."),
        ("[+] Zero Setup", "Packaged in Docker for simple 1-click startup on any Windows/Linux machine."),
        ("[+] Complete Privacy", "No data leaves the office network; 100% safe for official government documents."),
        ("[+] High Accuracy", "Direct page citations guarantee zero fake or made-up facts in summaries.")
    ]
    for title, desc in r_boxes:
        p = tf4_r.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{title}\n"
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.name = FONT_BODY
        r1.font.color.rgb = COLOR_ORANGE_HDR
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.name = FONT_BODY
        r2.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(10)

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide5, 5, "IMPACT AND BENEFITS")

    hdr_imp = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.3), Inches(6.0), Inches(0.4))
    hdr_imp.fill.solid()
    hdr_imp.fill.fore_color.rgb = COLOR_ORANGE_HDR
    hdr_imp.line.fill.background()
    p_hi = hdr_imp.text_frame.paragraphs[0]
    p_hi.text = "Impact"
    p_hi.font.size = Pt(13)
    p_hi.font.bold = True
    p_hi.font.name = FONT_BODY
    p_hi.font.color.rgb = COLOR_WHITE
    p_hi.alignment = PP_ALIGN.CENTER

    c_imp = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.7), Inches(6.0), Inches(2.4))
    c_imp.fill.solid()
    c_imp.fill.fore_color.rgb = COLOR_CARD_BG
    c_imp.line.color.rgb = COLOR_CARD_BORDER
    tf_imp = c_imp.text_frame
    tf_imp.word_wrap = True
    tf_imp.margin_left = Inches(0.15)
    tf_imp.margin_top = Inches(0.12)
    tf_imp.margin_right = Inches(0.15)

    imp_items = [
        ("Saves Analyst Time", "Reduces 4 to 6 hours of manual report writing to under 10 seconds."),
        ("Faster Decision Making", "Officers can quickly read 1-page executive memos during urgent meetings."),
        ("Regional Language Reach", "Converts central government policies into Hindi, Kannada, and local Indian languages.")
    ]
    for idx, (t, d) in enumerate(imp_items):
        p = tf_imp.add_paragraph() if idx > 0 else tf_imp.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"● {t}: "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = d
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(4)

    hdr_ben = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.3), Inches(6.0), Inches(0.4))
    hdr_ben.fill.solid()
    hdr_ben.fill.fore_color.rgb = COLOR_ORANGE_HDR
    hdr_ben.line.fill.background()
    p_hb = hdr_ben.text_frame.paragraphs[0]
    p_hb.text = "Benefits"
    p_hb.font.size = Pt(13)
    p_hb.font.bold = True
    p_hb.font.name = FONT_BODY
    p_hb.font.color.rgb = COLOR_WHITE
    p_hb.alignment = PP_ALIGN.CENTER

    c_ben = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.7), Inches(6.0), Inches(2.1))
    c_ben.fill.solid()
    c_ben.fill.fore_color.rgb = COLOR_CARD_BG
    c_ben.line.color.rgb = COLOR_CARD_BORDER
    tf_ben = c_ben.text_frame
    tf_ben.word_wrap = True
    tf_ben.margin_left = Inches(0.15)
    tf_ben.margin_top = Inches(0.12)
    tf_ben.margin_right = Inches(0.15)

    ben_items = [
        ("Cost Effective", "Eliminates high costs of outsourcing presentation design and document translation."),
        ("100% Data Privacy", "Sensitive government intelligence remains inside local servers without leaks."),
        ("Easy Mobile Access", "Field staff can listen to 60-second audio summaries on their mobile phones.")
    ]
    for idx, (t, d) in enumerate(ben_items):
        p = tf_ben.add_paragraph() if idx > 0 else tf_ben.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"● {t}: "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = COLOR_TEXT
        r2 = p.add_run()
        r2.text = d
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(4)

    hdr_story = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.3), Inches(5.8), Inches(0.4))
    hdr_story.fill.solid()
    hdr_story.fill.fore_color.rgb = COLOR_ORANGE_HDR
    hdr_story.line.fill.background()
    p_hs = hdr_story.text_frame.paragraphs[0]
    p_hs.text = "For Instance: User Story"
    p_hs.font.size = Pt(13)
    p_hs.font.bold = True
    p_hs.font.name = FONT_BODY
    p_hs.font.color.rgb = COLOR_WHITE
    p_hs.alignment = PP_ALIGN.CENTER

    c_story = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.7), Inches(5.8), Inches(5.1))
    c_story.fill.solid()
    c_story.fill.fore_color.rgb = COLOR_CARD_BG
    c_story.line.color.rgb = COLOR_CARD_BORDER
    tf_story = c_story.text_frame
    tf_story.word_wrap = True
    tf_story.margin_left = Inches(0.2)
    tf_story.margin_top = Inches(0.15)
    tf_story.margin_right = Inches(0.2)

    story_text = [
        ("Scenario:", "Rahul, an intelligence analyst at a government agency, receives a 50-page technical threat report at 9:00 AM. He has to present the key findings to senior officers in a 9:30 AM meeting."),
        ("Problem:", "Reading the full 50 pages and manually typing summary slides takes at least 3 hours. Finding the most critical points in a hurry can lead to missing important details."),
        ("Impact with OmniTransform:", "Rahul uploads the PDF to OmniTransform. In 10 seconds, the platform gives him:"),
        ("● 1-Page Summary Memo", "Clear bullet points with key highlights and numbers."),
        ("● 5 Presentation Slides", "Formatted and ready to project in the meeting room."),
        ("● Clickable Citations", "Rahul clicks any bullet point to show the exact page in the raw PDF."),
        ("Benefits:", "Rahul walks into the 9:30 AM meeting fully prepared with accurate, verified facts. What used to take 3 hours was finished in under 10 seconds.")
    ]

    for idx, (label, text) in enumerate(story_text):
        p = tf_story.add_paragraph() if idx > 0 else tf_story.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{label} "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = COLOR_NAVY if "Scenario" in label or "Problem" in label or "Impact" in label or "Benefits" in label else COLOR_TEXT
        
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(10.5)
        r2.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(4)

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES (With Dedicated Repo and Resource Specs)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header_footer(slide6, 6, "RESEARCH AND REFERENCES")

    tb6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.7), Inches(5.5))
    tf6 = tb6.text_frame
    tf6.word_wrap = True

    links_data = [
        ("MASTER RESOURCE REPOSITORY & CODE", DEDICATED_REPO_URL),
        ("OPEN-SOURCE DATASETS & MODELS GUIDE", f"{DEDICATED_REPO_URL}/blob/main/OPEN_SOURCE_RESOURCES_AND_DATASETS.md"),
        ("TECHNICAL DEVELOPMENT SPECIFICATION", f"{DEDICATED_REPO_URL}/blob/main/PROJECT_DEVELOPMENT_SPECIFICATION.md"),
        ("SMART INDIA HACKATHON 2026 OFFICIAL PORTAL (PS ID 26154)", "https://www.sih.gov.in/sih2026PS"),
        ("NATIONAL TECHNICAL RESEARCH ORGANISATION (NTRO) PORTAL", "https://ntro.gov.in"),
        ("RESEARCH ON RAG & CITATIONS (LEWIS ET AL. 2020)", "https://arxiv.org/abs/2005.11401")
    ]

    for idx, (title, url) in enumerate(links_data):
        p = tf6.add_paragraph() if idx > 0 else tf6.paragraphs[0]
        r1 = p.add_run()
        r1.text = "❖  "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_TEXT

        r2 = p.add_run()
        r2.text = f"{title}:  "
        r2.font.bold = True
        r2.font.size = Pt(13)
        r2.font.name = FONT_TITLE
        r2.font.color.rgb = COLOR_NAVY

        r3 = p.add_run()
        r3.text = url
        r3.font.underline = True
        r3.font.size = Pt(12)
        r3.font.name = FONT_BODY
        r3.font.color.rgb = COLOR_FOOTER_BLUE
        p.space_after = Pt(10)

    prs.save(output_path)
    print(f"Deck with dedicated repo {DEDICATED_REPO_URL} created successfully at: {output_path}")

if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "SIH2026_Idea_Presentation_PS26154.pptx"
    build_exact_sih_deck(out)
